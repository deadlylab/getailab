#!/usr/bin/env python3
"""
CryptO'Brien Pty Ltd - Lab Agent V2 (The Research Sandbox)
Port: 5035 | Division: R&D (GetAiLab)
Role: High-Performance Execution & Artifact Storage
"""
import os
import sys
import re
import json
import time
import subprocess
import sqlite3
import asyncio
import shutil
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional
from flask import Flask, request, jsonify, send_from_directory, redirect, make_response
import random
import math
import hmac
import hashlib
import base64
from collections import defaultdict
from flask_cors import CORS
from playwright.async_api import async_playwright
from llm.adapter import create_default_adapter, get_env_provider_config
from llm.sauron_core import extract_with_adapter
import requests
from bs4 import BeautifulSoup
from markdownify import markdownify as md

# ==========================================
# GETAILAB LIBRARY — data/labs/<lab_id>/ scientist books + codex
# ==========================================
GETAILAB_LIB = None
GETAILAB_LIBRARY_ENABLED = False
GETAILAB_INTEGRITY_ENABLED = False
GETAILAB_LEARNING_ENABLED = False
GETAILAB_TICKETS_ENABLED = False
get_loop_ticket_tracker = None
try:
    from getailab.tickets import get_loop_ticket_tracker as _get_loop_ticket_tracker
    get_loop_ticket_tracker = _get_loop_ticket_tracker
    GETAILAB_TICKETS_ENABLED = True
except Exception as _ticket_e:
    print(f"[JobTickets] Integration unavailable ({_ticket_e}).")

try:
    _project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    if _project_root not in sys.path:
        sys.path.insert(0, _project_root)
    from getailab.library.service import (
        add_scientist_reference,
        get_library,
        get_scientist_references,
        reindex_library,
        valid_scientist_name,
    )
    _active_lab = os.getenv("LAB_ID", "example").strip() or "example"
    GETAILAB_LIB = get_library(lab_id=_active_lab)
    GETAILAB_LIBRARY_ENABLED = True
except Exception as _lib_e:
    add_scientist_reference = None
    get_scientist_references = None
    reindex_library = None
    valid_scientist_name = None
    print(f"[GetAiLabLibrary] Dashboard integration unavailable ({_lib_e}). Using legacy loop views.")

try:
    from getailab.integrity.verify import (
        crush_test_indexes,
        crush_test_vault,
        merkle_status,
        scan_integrity_targets,
        verify_full,
    )
    from getailab.integrity.signing import (
        attest_vault,
        generate_keypair,
        sign_merkle_tree,
        signing_available,
        signing_status,
        verify_merkle_signature,
    )
    GETAILAB_INTEGRITY_ENABLED = True
except Exception as _integrity_e:
    crush_test_indexes = None
    crush_test_vault = None
    merkle_status = None
    scan_integrity_targets = None
    verify_full = None
    attest_vault = None
    generate_keypair = None
    sign_merkle_tree = None
    signing_available = None
    signing_status = None
    verify_merkle_signature = None
    print(f"[Integrity] API unavailable ({_integrity_e}).")

try:
    from getailab.learning import get_adaptive_learner
    from getailab.gabby.gabby import Gabby
    GETAILAB_LEARNING_ENABLED = True
except Exception as _learning_e:
    get_adaptive_learner = None
    Gabby = None
    print(f"[AdaptiveLearner] Integration unavailable ({_learning_e}).")

try:
    from fpdf import FPDF
except ImportError:
    FPDF = None
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    Presentation = None

app = Flask(__name__)
CORS(app)

try:
    from dotenv import load_dotenv
    load_dotenv(os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), ".env"))
except ImportError:
    pass

AGENT_PORT = int(os.getenv('LAB_PORT', '5035'))
ORACLE_URL = os.getenv("ORACLE_URL", "http://localhost:5024").rstrip("/")
CHAT_ADAPTER = None
# Pathing for the permanent artifacts directory
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DASHBOARD_DIR = os.path.join(os.path.dirname(BASE_DIR), 'dashboard', 'frontend')
ROOT_DIR = os.path.dirname(BASE_DIR)  # for loading user notes / loop reports

# ---------------------------------------------------------------------------
# Dashboard gate — hide the frontend wankfest until product wiring is ready.
# Set GETAILAB_DASHBOARD_PASSWORD in .env to enable. Empty = open (local only).
#
# Applies to EVERY lab: chimera, university, example, old_mate, rf_research, …
# all boot the same lab/app_lab.py. Machine dialectic routes stay open so a
# dashboard password never 401's /execute or /literature (product-lab Loop 1).
# Restart the lab process after changing this module.
# ---------------------------------------------------------------------------
DASHBOARD_PASSWORD = (os.getenv("GETAILAB_DASHBOARD_PASSWORD") or "").strip()
DASHBOARD_USER = (os.getenv("GETAILAB_DASHBOARD_USER") or "getailab").strip()
DASHBOARD_GATE_SECRET = (
    os.getenv("GETAILAB_DASHBOARD_GATE_SECRET")
    or os.getenv("GETAILAB_DASHBOARD_PASSWORD")
    or "getailab-dev-gate"
).strip()
DASHBOARD_COOKIE = "getailab_gate"
# HTML UI only — dialectic machine endpoints stay open so run_lab / scientists
# are not 401'd by the public dashboard password (Loop 1 university failure mode).
_GATE_OPEN_PATHS = {
    "/health",
    "/gate",
    "/gate/login",
    "/gate/logout",
    "/favicon.ico",
}
# Prefixes used by run_lab.py + agents — never require browser cookie/basic.
# Keep in sync with create_lab / run_*.sh comments; shared by all lab_ids.
_GATE_OPEN_PREFIXES = (
    "/execute",
    "/develop/",
    "/literature/",
    "/vision/",
    "/web/",
)


def _gate_enabled() -> bool:
    return bool(DASHBOARD_PASSWORD)


def _gate_token() -> str:
    raw = f"{DASHBOARD_USER}:{DASHBOARD_PASSWORD}".encode()
    dig = hmac.new(DASHBOARD_GATE_SECRET.encode(), raw, hashlib.sha256).digest()
    return base64.urlsafe_b64encode(dig).decode().rstrip("=")


def _gate_ok() -> bool:
    if not _gate_enabled():
        return True
    # HTTP Basic (easy for curl / API clients)
    auth = request.authorization
    if auth and auth.username == DASHBOARD_USER and auth.password == DASHBOARD_PASSWORD:
        return True
    # Cookie after form login
    cookie = request.cookies.get(DASHBOARD_COOKIE, "")
    if cookie and hmac.compare_digest(cookie, _gate_token()):
        return True
    return False


def _gate_landing_html(error: str = "") -> str:
    err = (
        f'<p style="color:#f87171;margin:0 0 1rem;font-size:0.9rem">{error}</p>'
        if error
        else ""
    )
    return f"""<!doctype html>
<html lang="en"><head>
<meta charset="utf-8"/><meta name="viewport" content="width=device-width,initial-scale=1"/>
<title>GetAiLab — Authorized Access</title>
<style>
  html,body{{margin:0;min-height:100%;background:#0a0a12;color:#e0e7ff;
    font-family:Inter,system-ui,sans-serif;display:flex;align-items:center;justify-content:center}}
  .card{{width:min(420px,92vw);padding:2rem;border:1px solid #1e1e2e;border-radius:16px;
    background:linear-gradient(160deg,#11111a,#0a0a12);box-shadow:0 0 60px rgba(192,132,252,.08)}}
  h1{{font-family:Georgia,serif;font-style:italic;font-weight:500;font-size:1.5rem;margin:0 0 .5rem;color:#e8e4dc}}
  p.sub{{color:#94a3b8;font-size:.9rem;line-height:1.5;margin:0 0 1.5rem}}
  label{{display:block;font-size:.75rem;color:#64748b;margin:0 0 .35rem;letter-spacing:.04em;text-transform:uppercase}}
  input{{width:100%;box-sizing:border-box;padding:.75rem .9rem;border-radius:10px;border:1px solid #2a2a3d;
    background:#0a0a12;color:#e0e7ff;margin-bottom:1rem;font-size:1rem}}
  input:focus{{outline:none;border-color:#c084fc}}
  button{{width:100%;padding:.8rem;border:none;border-radius:10px;background:#fcd34d;color:#0a0a12;
    font-weight:600;font-size:.95rem;cursor:pointer}}
  button:hover{{filter:brightness(1.05)}}
  .tag{{display:inline-block;font-size:.65rem;letter-spacing:.12em;text-transform:uppercase;
    color:#67e8f9;border:1px solid #164e63;padding:.2rem .5rem;border-radius:999px;margin-bottom:1rem}}
</style></head><body>
<div class="card">
  <div class="tag">Private · Contract preview</div>
  <h1>GetAiLab</h1>
  <p class="sub">Public UI is gated while we finish product wiring.
  Authorized operators only — the research OS is not the landing page wankfest.</p>
  {err}
  <form method="post" action="/gate/login">
    <label for="user">User</label>
    <input id="user" name="user" autocomplete="username" value="{DASHBOARD_USER}"/>
    <label for="password">Password</label>
    <input id="password" name="password" type="password" autocomplete="current-password" required autofocus/>
    <button type="submit">Enter lab</button>
  </form>
</div>
</body></html>"""


@app.before_request
def _dashboard_gate():
    if not _gate_enabled():
        return None
    path = request.path or "/"
    if path in _GATE_OPEN_PATHS or path.startswith("/gate/"):
        return None
    # Dialectic OS endpoints — open (sandbox, lit, Sauron). Gate is for HTML UI only.
    if any(path == p or path.startswith(p) for p in _GATE_OPEN_PREFIXES):
        return None
    # Authenticated operator — full access including /api/* dashboard JSON
    if _gate_ok():
        return None
    # Unauthenticated: JSON APIs get 401 JSON; pages get the login HTML
    if path.startswith("/api/"):
        return jsonify({
            "error": "dashboard_gate",
            "message": "Set GETAILAB_DASHBOARD_PASSWORD cookie or HTTP Basic auth.",
        }), 401
    return make_response(_gate_landing_html(), 401)


@app.route("/gate", methods=["GET"])
@app.route("/gate/login", methods=["GET", "POST"])
def gate_login():
    if not _gate_enabled():
        return redirect("/")
    if request.method == "GET":
        if _gate_ok():
            return redirect("/")
        return make_response(_gate_landing_html(), 200)
    user = (request.form.get("user") or "").strip()
    password = request.form.get("password") or ""
    if user == DASHBOARD_USER and password == DASHBOARD_PASSWORD:
        resp = make_response(redirect("/"))
        resp.set_cookie(
            DASHBOARD_COOKIE,
            _gate_token(),
            httponly=True,
            samesite="Lax",
            max_age=60 * 60 * 24 * 14,  # 14 days
        )
        return resp
    return make_response(_gate_landing_html("Wrong password. Try again."), 401)


@app.route("/gate/logout", methods=["GET", "POST"])
def gate_logout():
    resp = make_response(redirect("/gate"))
    resp.set_cookie(DASHBOARD_COOKIE, "", expires=0)
    return resp



def _resolve_lab_paths() -> tuple:
    """Per-lab sandbox DB + artifacts — the example lab uses lab/; forged labs use data/labs/<id>/."""
    try:
        from getailab.lab_config import lab_artifacts_dir, lab_results_db_path, get_lab_id
        lid = get_lab_id()
        return str(lab_results_db_path(lid)), str(lab_artifacts_dir(lid))
    except Exception:
        return (
            os.path.join(BASE_DIR, 'lab_results.db'),
            os.path.join(BASE_DIR, 'artifacts'),
        )


DB_PATH, ARTIFACTS_DIR = _resolve_lab_paths()

try:
    from getailab.lab_config import get_lab_id, is_chimera_clone_lab, is_chimera_lab, lab_reports_dir, resolve_lab_paths
    ACTIVE_LAB_ID = get_lab_id()
    LAB_PATHS = resolve_lab_paths(ACTIVE_LAB_ID)
except Exception:
    ACTIVE_LAB_ID = os.getenv("LAB_ID", "example").strip() or "example"
    _is_clone = ACTIVE_LAB_ID in ("chimera_clone", "chimera")
    LAB_PATHS = {
        "lab_id": ACTIVE_LAB_ID,
        "artifacts": ARTIFACTS_DIR,
        "results_db": DB_PATH,
        "is_chimera_clone": _is_clone,
        "is_chimera": _is_clone,
    }
    is_chimera_clone_lab = lambda lab_id=None: (lab_id or ACTIVE_LAB_ID) in ("chimera_clone", "chimera")  # type: ignore
    is_chimera_lab = is_chimera_clone_lab  # type: ignore
    lab_reports_dir = lambda lab_id=None: ROOT_DIR  # type: ignore

# Lab DB alias fix (existing _query_lab_db referenced undefined LAB_DB)
LAB_DB = DB_PATH

# Ensure the lab has a physical workspace for data
os.makedirs(ARTIFACTS_DIR, exist_ok=True)
os.makedirs(DASHBOARD_DIR, exist_ok=True)


def _lab_ops_log_path() -> str:
    """Durable ops log for implement/execute (not Flask access logs)."""
    vault = (os.getenv("GETAILAB_LAB_ROOT") or "").strip()
    if vault and os.path.isdir(vault):
        log_dir = os.path.join(vault, "logs")
    else:
        log_dir = os.path.join(os.path.dirname(ARTIFACTS_DIR), "logs")
    os.makedirs(log_dir, exist_ok=True)
    return os.path.join(log_dir, "lab_ops.jsonl")


def _append_lab_ops(event: dict) -> None:
    """Append one JSON line to lab_ops.jsonl; never raise into request path."""
    try:
        payload = dict(event or {})
        payload.setdefault("ts", datetime.utcnow().isoformat() + "Z")
        payload.setdefault("lab_id", ACTIVE_LAB_ID)
        with open(_lab_ops_log_path(), "a", encoding="utf-8") as fh:
            fh.write(json.dumps(payload, default=str) + "\n")
    except Exception:
        pass


def _evaluate_execute_success(returncode: int, stdout: str) -> tuple:
    """Honor RESULT PASS/FAIL lines — exit 0 alone is not enough.

    Returns (success: bool, verdict: str).
    Match only line-leading RESULT tokens (avoid 'no result line' false positives).
    """
    text = stdout or ""
    last_fail_idx = -1
    last_pass_idx = -1
    for i, line in enumerate(text.splitlines()):
        s = line.strip().upper()
        if re.match(r"^RESULT\s*:\s*FAIL\b", s) or re.match(r"^RESULT\s+FAIL\b", s):
            last_fail_idx = i
        elif re.match(r"^RESULT\s*:\s*PASS\b", s) or re.match(r"^RESULT\s+PASS\b", s):
            last_pass_idx = i

    if last_fail_idx >= 0 and last_fail_idx >= last_pass_idx:
        return False, "result_fail"

    if returncode != 0:
        return False, "nonzero_exit"

    if last_pass_idx >= 0:
        return True, "result_pass"

    # No clear RESULT line: keep process exit as authority (prompt now requires
    # RESULT PASS/FAIL + sys.exit). Ops log records missing_result for coaching.
    return True, "missing_result_line"

def init_db():
    try:
        from getailab.lab_config import ensure_lab_results_db
        ensure_lab_results_db()
    except Exception:
        conn = sqlite3.connect(DB_PATH)
        conn.execute("""
            CREATE TABLE IF NOT EXISTS lab_experiments (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                loop_id TEXT,
                agent_name TEXT,
                experiment_name TEXT,
                code TEXT,
                stdout TEXT,
                stderr TEXT,
                success BOOLEAN,
                execution_time_ms INTEGER,
                artifacts_json TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        conn.commit()
        conn.close()

class SauronVision:
    """Lab /vision/extract endpoint — fast HTTP+text by default (matches sauron_vision.py)."""

    def __init__(self):
        self.adapter = create_default_adapter()

    async def extract(self, url, user_query):
        import os
        import requests as _req

        use_browser = os.getenv("SAURON_BROWSER", "").lower() in ("1", "true", "yes", "on")
        use_vision = os.getenv("SAURON_VISION", "").lower() in ("1", "true", "yes", "on")
        text_chars = int(os.getenv("SAURON_TEXT_CHARS", "6000"))
        markdown_text = ""
        screenshot_bytes = None

        if use_browser:
            async with async_playwright() as p:
                browser = await p.chromium.launch(
                    headless=True, args=["--no-sandbox", "--disable-dev-shm-usage"]
                )
                context = await browser.new_context(viewport={"width": 1280, "height": 800})
                page = await context.new_page()
                try:
                    await page.goto(url, wait_until="domcontentloaded", timeout=20000)
                    try:
                        await page.wait_for_load_state("networkidle", timeout=5000)
                    except Exception:
                        pass
                    content = await page.content()
                    soup = BeautifulSoup(content, "lxml")
                    for tag in soup(["script", "style", "nav", "footer", "noscript"]):
                        tag.decompose()
                    markdown_text = md(str(soup))
                    if use_vision:
                        screenshot_bytes = await page.screenshot(full_page=False, type="jpeg", quality=60)
                finally:
                    await browser.close()
        else:
            try:
                resp = _req.get(
                    url,
                    headers={"User-Agent": "Mozilla/5.0 (compatible; GetAiLab-Sauron/1.0)"},
                    timeout=int(os.getenv("SAURON_HTTP_TIMEOUT", "15")),
                )
                resp.raise_for_status()
                soup = BeautifulSoup(resp.text, "lxml")
                for tag in soup(["script", "style", "nav", "footer", "noscript", "iframe"]):
                    tag.decompose()
                markdown_text = md(str(soup))
            except Exception as e:
                return json.dumps({
                    "error": "capture_failed",
                    "hint": f"HTTP fetch failed ({e}). Set SAURON_BROWSER=1 for JS-heavy pages.",
                    "llm": self.adapter.get_info(),
                })

        data = extract_with_adapter(
            self.adapter,
            user_query,
            image_bytes=screenshot_bytes if use_vision else None,
            text_context=markdown_text[:text_chars] if markdown_text else None,
            max_text_chars=text_chars,
            fast=not use_vision,
        )
        if data:
            return json.dumps(data)
        return json.dumps({
            "error": "extraction_failed",
            "hint": "Check LLM_PROVIDER and that your backend is running (ollama serve, API key, etc.)",
            "llm": self.adapter.get_info(),
        })


# ============================================
# GET AI LAB DASHBOARD BACKEND SAUCE
# Real data, live pulses, inspiration engine
# ============================================

_AGENT_COLOR_PALETTE = [
    '#fcd34d', '#a5b4fc', '#c084fc', '#67e8f9', '#f472b6',
    '#fbbf24', '#a78bfa', '#34d399', '#60a5fa', '#fb7185', '#f97316',
]

_AGENT_PERSONAS_DEFAULT = {
    'albert': {'role': 'Spacetime Geometry', 'color': '#fcd34d'},
    'bohr': {'role': 'Complementarity & Limits', 'color': '#a5b4fc'},
    'heisenberg': {'role': 'Matrix Mechanics', 'color': '#c084fc'},
    'alan': {'role': 'Epistemology & Belief', 'color': '#67e8f9'},
    'brian': {'role': 'Cognitive Higgs Fields', 'color': '#f472b6'},
    'carl': {'role': 'Archetypes & Collapse', 'color': '#fbbf24'},
    'neil': {'role': 'Cosmic Emergence', 'color': '#a78bfa'},
    'roger': {'role': 'Orch-OR & Microtubules', 'color': '#34d399'},
    'emmy': {'role': 'Symmetry & Topology', 'color': '#60a5fa'},
    'andrew': {'role': 'Neuro-Metabolic Protocols', 'color': '#fb7185'},
    'tesla': {'role': 'Resonance & Coupled Systems', 'color': '#f97316'},
}


def _load_agent_personas() -> Dict[str, Dict[str, str]]:
    """Sync dashboard/chat personas with active squad YAML."""
    try:
        from personas.loader import get_squad_names, get_persona
        out: Dict[str, Dict[str, str]] = {}
        idx = 0
        for name in get_squad_names():
            if name == 'oracle':
                continue
            try:
                p = get_persona(name)
                role = p.get('display_role') or p.get('role') or name.title()
            except Exception:
                role = _AGENT_PERSONAS_DEFAULT.get(name, {}).get('role', name.title())
            color = (
                _AGENT_PERSONAS_DEFAULT.get(name, {}).get('color')
                or _AGENT_COLOR_PALETTE[idx % len(_AGENT_COLOR_PALETTE)]
            )
            out[name] = {'role': role, 'color': color}
            idx += 1
        return out if out else dict(_AGENT_PERSONAS_DEFAULT)
    except Exception:
        return dict(_AGENT_PERSONAS_DEFAULT)


AGENT_PERSONAS = _load_agent_personas()

# Legacy local reference-lab curated voices — not used for example or forged labs
_REFERENCE_LAB_CURATED_QUOTES: Dict[str, str] = {
    "albert": (
        "Bohr, my dear friend... 'God does not play dice with the universe — and neither should our agents.' "
        "The geometry must be felt. Show me the manifold or we describe shadows."
    ),
    "bohr": (
        "Albert, my dear friend, you would argue there must be a hidden variable... I must respectfully disagree. "
        "'Albert, stop telling God what to do!' Complementarity is the only honest description."
    ),
    "heisenberg": (
        "This is poetry, not physics. We need the operators and their commutation relations. "
        "[x, p] = iħ is carved into existence. Vague notions of 'dialectic' are insufficient."
    ),
    "alan": (
        "True intelligence does not collapse the wavefunction too early. "
        "We are building the conditions under which knowing can occur."
    ),
    "brian": (
        "What gives the 'information' its mass? Show me the Lagrangian. What is the coupling constant?"
    ),
    "carl": (
        "Repressing the Shadow ensures it erupts as hallucination. Tokens are symbols, not statistics."
    ),
    "neil": (
        "If it is not scale-invariant, it is probably not fundamental. "
        "The carbon in our neurons was forged in the heart of a dying star."
    ),
    "roger": (
        "The brain must be doing something non-algorithmic. "
        "We are taking the ashes of dead stars and coaxing them to dream."
    ),
    "emmy": (
        "This is not a homomorphism. Show me the symmetry. What is conserved? Formalize or fail."
    ),
    "andrew": (
        "The brain runs on a sandwich and a glass of water. Do not bend the biology."
    ),
}


def _refresh_agent_personas() -> Dict[str, Dict[str, str]]:
    """Reload squad from active PERSONAS_YAML (per-lab isolation)."""
    global AGENT_PERSONAS
    AGENT_PERSONAS = _load_agent_personas()
    return AGENT_PERSONAS


def _extract_quote_from_persona(name: str, persona: Dict[str, Any]) -> str:
    """Pull a council voice line from this lab's persona YAML only."""
    sp = (persona.get("system_prompt") or "").strip()
    in_examples = False
    for line in sp.splitlines():
        stripped = line.strip()
        if re.search(r"example interaction", stripped, re.I):
            in_examples = True
            continue
        if in_examples and stripped.startswith("- "):
            quote = stripped[2:].strip().strip('"').strip("'")
            if len(quote) > 35:
                return quote[:340]
        if in_examples and stripped.startswith("**") and "example" not in stripped.lower():
            in_examples = False
    for line in sp.splitlines():
        s = line.strip()
        if s.startswith(("- ", '- "', "- '")) and len(s) > 45:
            return s.lstrip("- ").strip('"\'')[:340]
    expertise = persona.get("expertise") or []
    if expertise:
        e0 = expertise[0] if isinstance(expertise[0], str) else str(expertise[0])
        if len(e0) > 30:
            return e0[:320]
    focus = (
        persona.get("implement_focus")
        or persona.get("display_role")
        or persona.get("role")
        or name.replace("_", " ").title()
    )
    return (
        f"{name.replace('_', ' ').title()} — {focus}. "
        "Demand testable hypotheses and auditable artifacts for this lab's agenda."
    )[:340]


def _get_lab_quotes() -> List[tuple]:
    """Quotes restricted to the active lab's squad — no the example lab bleed on forged labs."""
    from personas.loader import get_persona

    _refresh_agent_personas()
    quotes: List[tuple] = []
    use_curated = is_chimera_lab(ACTIVE_LAB_ID)
    for name in AGENT_PERSONAS:
        if use_curated and name in _REFERENCE_LAB_CURATED_QUOTES:
            quotes.append((name, _REFERENCE_LAB_CURATED_QUOTES[name]))
            continue
        try:
            p = get_persona(name)
            quotes.append((name, _extract_quote_from_persona(name, p)))
        except Exception:
            info = AGENT_PERSONAS.get(name, {})
            quotes.append((
                name,
                f"{name.replace('_', ' ').title()} — {info.get('role', 'Research')}: "
                "argue with data, charts, and reproducible artifacts.",
            ))
    if not quotes:
        quotes.append(("oracle", "This lab's council is ready — forge the first loop when you are."))
    return quotes


def _lab_research_teasers() -> List[str]:
    teasers = list(RESEARCH_TEASERS)
    try:
        from getailab.lab_config import load_lab_config
        agenda = (load_lab_config(ACTIVE_LAB_ID).get("research_agenda") or "").strip()
        if agenda:
            teasers.extend([
                f"What experiment best advances: {agenda[:140]}?",
                f"Which artifact in {ACTIVE_LAB_ID} should inform the next hypothesis?",
                f"How do we chart and falsify claims about {agenda[:100]}?",
            ])
    except Exception:
        pass
    return teasers

RESEARCH_TEASERS = [
    "What experiment would best test whether multi-agent review improves hypothesis quality over a single model?",
    "How should prior loop artifacts be indexed so the squad can build on them without repeating work?",
    "What metrics distinguish a useful synthesis from a polished summary of existing ideas?",
    "When does sandbox execution fail to validate a hypothesis, and what should the loop do next?",
    "How can complementarity between formal and intuitive reasoning be made explicit in experiment design?",
    "What are the trade-offs between breadth and depth when assigning problems to a specialist squad?",
    "How do we detect when a research loop has converged versus when it needs a new problem vector?",
    "What makes an artifact auditable enough to survive independent review?",
    "How should conflicting hypotheses from different scientists be preserved in the record?",
    "What is the minimum viable problem statement for a productive first loop?",
]

def _agora_db_path() -> str:
    try:
        from getailab.lab_config import agora_db_path as _agora_path
        return str(_agora_path(ACTIVE_LAB_ID))
    except Exception:
        return str(LAB_PATHS.get("agora_db", os.path.join(ROOT_DIR, "chimera_lab.db")))


def _query_agora_db(query, params=()):
    """Query the active lab's loop DB (chimera_lab.db or data/labs/<id>/agora.db)."""
    try:
        path = _agora_db_path()
        if not os.path.isfile(path):
            return []
        conn = sqlite3.connect(path, timeout=5)
        rows = conn.execute(query, params).fetchall()
        conn.close()
        return rows
    except Exception:
        return []


def _query_chimera_db(query, params=()):
    """Backward-compatible alias."""
    return _query_agora_db(query, params)

def _query_lab_db(query, params=()):
    try:
        conn = sqlite3.connect(DB_PATH, timeout=5)
        rows = conn.execute(query, params).fetchall()
        conn.close()
        return rows
    except Exception:
        return []

def compute_real_stats():
    """Pull authentic numbers from the twin databases + filesystem for the living dashboard."""
    loops_rows = _query_chimera_db("SELECT COUNT(*), MAX(loop_id) FROM agora_loops")
    loops_count = loops_rows[0][0] if loops_rows else 0
    max_loop = loops_rows[0][1] if loops_rows and loops_rows[0][1] else 0

    # Experiments and artifacts from lab
    exp_rows = _query_lab_db("SELECT COUNT(*), agent_name FROM lab_experiments GROUP BY agent_name")
    total_experiments = sum(r[0] for r in exp_rows)
    agent_contribs = {r[1]: r[0] for r in exp_rows}

    # Total artifacts (parse JSON length)
    art_rows = _query_lab_db("SELECT artifacts_json FROM lab_experiments")
    total_artifacts = 0
    for (aj,) in art_rows:
        try:
            if aj:
                arr = json.loads(aj)
                total_artifacts += len(arr) if isinstance(arr, list) else 0
        except:
            pass

    # Knowledge growth data: artifacts per loop (approximate via experiments + real files)
    growth = []
    loop_art = defaultdict(int)
    for (lid, aj) in _query_lab_db("SELECT loop_id, artifacts_json FROM lab_experiments"):
        try:
            if aj:
                cnt = len(json.loads(aj)) if aj else 0
                loop_art[str(lid)] += cnt
        except:
            pass
    # Sort by numeric loop
    sorted_loops = sorted(loop_art.items(), key=lambda x: int(x[0]) if x[0].isdigit() else 0)
    cumulative = 0
    for lid, cnt in sorted_loops[-8:]:  # last 8 for chart sanity
        cumulative += cnt
        growth.append({"loop": int(lid) if lid.isdigit() else lid, "artifacts": cnt, "cumulative": cumulative})

    _refresh_agent_personas()
    squad_size = max(1, len(AGENT_PERSONAS))

    # "Books" = full consensus artefacts written
    books = _query_chimera_db("SELECT COUNT(*) FROM agora_loops WHERE consensus_artefact IS NOT NULL")[0][0] if _query_chimera_db("SELECT COUNT(*) FROM agora_loops WHERE consensus_artefact IS NOT NULL") else loops_count // 2 + 2
    # GETAILABLIBRARY real count: every loop = multiple pages in the book
    library_pages = 0
    library_insp = 88
    recent_lib = []
    if GETAILAB_LIBRARY_ENABLED and GETAILAB_LIB:
        try:
            b = GETAILAB_LIB.get_or_create_default_book()
            library_pages = len(b.page_ids)
            books = max(books, library_pages // 4 + 1)
            # Pull live summary for UI charts + personal sauce
            lib_sum = GETAILAB_LIB.get_recent_library_summary(7)
            library_insp = lib_sum.get("inspiration", library_insp)
            recent_lib = lib_sum.get("recent_pages", [])[:5]
        except Exception:
            pass

    last_loop_rows = _query_chimera_db("SELECT start_time FROM agora_loops ORDER BY loop_id DESC LIMIT 1")

    from getailab.resonance import (
        build_trajectory,
        compute_inspiration_score,
        compute_resonance_streak,
        compute_synthesis_rate,
    )

    stats = {
        "loops_completed": loops_count,
        "max_loop_id": max_loop or 0,
        "total_artifacts": total_artifacts,
        "total_experiments": total_experiments,
        "books_created": books,
        "agent_contributions": agent_contribs,
        "knowledge_growth": growth,
        "last_loop_time": last_loop_rows[0][0] if last_loop_rows else None,
        "library_pages": library_pages,
        "library_resonance": library_insp,
        "recent_library_activity": recent_lib,
        "lab_id": ACTIVE_LAB_ID,
        "squad_size": squad_size,
    }
    stats["research_progress"] = get_research_progress(stats, squad_size)
    stats["inspiration_score"] = compute_inspiration_score(stats, ACTIVE_LAB_ID, squad_size)
    stats["resonance_streak"] = compute_resonance_streak(ACTIVE_LAB_ID)
    stats["synthesis"] = compute_synthesis_rate(ACTIVE_LAB_ID)
    stats["resonance_trajectory"] = build_trajectory(stats, ACTIVE_LAB_ID, squad_size)
    dirs = load_user_directives(1)
    stats["recent_directive"] = dirs[0]["note"] if dirs else "Run a loop to seed the lab record."
    stats["recent_directive_loop"] = dirs[0]["loop"] if dirs else 0
    if stats.get("library_pages"):
        stats["research_progress"] = min(99, stats["research_progress"] + min(9, stats["library_pages"] // 12))
    return stats

def compute_reminders(stats):
    """Smart, contextual, tyre-pumping reminders derived from real state.
    NOW TIED TO YOUR ACTUAL NOTES & DIRECTIVES from the library of loops + GetAiLabLibrary pages."""
    reminders = []
    last_id = stats.get("max_loop_id", 16)
    directives = load_user_directives(4)
    lib_pages = stats.get("library_pages", 0)

    reminders.append({
        "id": "r1",
        "type": "review",
        "title": f"Loop #{last_id} synthesis is resting in the Agora.",
        "body": "The Council wove a rich Consensus Artefact. Your eye on the geodesics could reveal the next invariant.",
        "action": "Review Synthesis",
        "boost": 4
    })
    if lib_pages > 30 or stats.get("total_artifacts", 0) > 700:
        new_pages = max(3, lib_pages // 12) if lib_pages else random.randint(7,14)
        reminders.append({
            "id": "r2",
            "type": "library",
            "title": f"{new_pages} new pages crystallized in the GetAiLabLibrary Vault since last resonance.",
            "body": "Plots, matrices, belief-evolution signatures, and your own architect notes await pattern recognition. The codex grows because of you.",
            "action": "Search Codex",
            "boost": 3
        })
    reminders.append({
        "id": "r3",
        "type": "continue",
        "title": "Ready for a follow-up loop.",
        "body": f"Prior work is in the record. {lib_pages} library pages indexed. Consider a follow-up problem that tests the last synthesis.",
        "action": "Start Next Loop",
        "boost": 5
    })
    if stats.get("inspiration_score", 80) > 85 or lib_pages > 60:
        reminders.append({
            "id": "r4",
            "type": "celebrate",
            "title": "Strong research activity.",
            "body": f"{stats['loops_completed']} loops. {stats['total_artifacts']} artifacts. {lib_pages or 0} library pages.",
            "action": "View Stats",
            "boost": 3
        })

    # NEW: Smart directive-tied reminder. References YOUR notes directly — deeply personal sauce.
    if directives:
        d = directives[0]
        reminders.append({
            "id": "r5",
            "type": "directive",
            "title": f"Revisit your directive from Loop #{d.get('loop', 12)}.",
            "body": d["note"][:170] + ("..." if len(d["note"])>170 else ""),
            "action": "Build on This Note",
            "boost": 6,
            "directive_ref": d
        })
    return reminders[:5]

def get_random_inspiration(agent_filter: Optional[str] = None) -> Dict[str, Any]:
    quotes = _get_lab_quotes()
    if agent_filter:
        filt = agent_filter.lower().strip()
        matched = [(a, q) for a, q in quotes if a.lower() == filt]
        if matched:
            quotes = matched
    agent, quote = random.choice(quotes)
    persona = AGENT_PERSONAS.get(agent, {})
    display = agent.replace("_", " ").title()
    return {
        "agent": display,
        "agent_key": agent,
        "role": persona.get("role", "Council Voice"),
        "quote": quote,
        "teaser": random.choice(_lab_research_teasers()),
        "color": persona.get("color", "#a5b4fc"),
        "lab_id": ACTIVE_LAB_ID,
    }


def _chat_adapter():
    global CHAT_ADAPTER
    if CHAT_ADAPTER is None:
        CHAT_ADAPTER = create_default_adapter()
    return CHAT_ADAPTER


def _pick_chat_agent(user_msg: str) -> str:
    """Route to a named scientist when explicitly addressed; else Oracle."""
    msg = (user_msg or "").lower()
    for name in AGENT_PERSONAS:
        if re.search(rf"\b{re.escape(name)}\b", msg):
            return name
    return "oracle"


def _format_chat_history(history: list) -> str:
    lines = []
    for turn in (history or [])[-4:]:
        if isinstance(turn, dict):
            user = turn.get("user") or turn.get("message") or ""
            agent = turn.get("agent") or "ORACLE"
            reply = turn.get("reply") or ""
            if user:
                lines.append(f"User: {user}")
            if reply:
                lines.append(f"{agent}: {reply[:500]}")
    return "\n".join(lines)


def _generate_council_reply(user_msg: str, history: list, agent_key: str) -> str:
    _refresh_agent_personas()
    persona_meta = AGENT_PERSONAS.get(agent_key, {"role": "Council Orchestrator"})
    role = persona_meta.get("role", "Council")
    agent_label = agent_key.replace("_", " ").title() if agent_key != "oracle" else "ORACLE"

    yaml_prompt = ""
    try:
        from personas.loader import get_persona
        yaml_prompt = (get_persona(agent_key).get("system_prompt") or "")[:6000]
    except Exception:
        pass

    if yaml_prompt:
        system = (
            yaml_prompt
            + f"\n\nLIVE COUNCIL CHAT — lab '{ACTIVE_LAB_ID}' only. "
            "Answer the user's message directly and usefully. "
            "Do not reference other labs or scientists outside this squad. "
            "Do not invent loop results — if unsure, say what experiment would test it."
        )
    else:
        system = f"""You are {agent_label} in GetAiLab lab '{ACTIVE_LAB_ID}'.
Role: {role}.
Answer the researcher's message directly, clearly, and usefully.
Stay within this lab's research domain. Do not cite other departments' work.
Give complete answers — finish your thought; do not trail off mid-sentence.
Do not invent loop results — if unsure, say what would need a dialectic loop to test."""

    hist = _format_chat_history(history)
    prompt_parts = []
    if hist:
        prompt_parts.append("RECENT CONVERSATION:\n" + hist)
    prompt_parts.append("USER MESSAGE:\n" + user_msg)
    prompt_parts.append("Reply as " + agent_label + ":")

    raw = _chat_adapter().generate(
        prompt="\n\n".join(prompt_parts),
        system_prompt=system,
    )
    text = (raw or "").strip()
    if text.startswith("ERROR:"):
        raise RuntimeError(text)
    max_chars = int(os.getenv("CHAT_MAX_REPLY_CHARS", "8000"))
    if len(text) > max_chars:
        cut = text[:max_chars].rsplit(" ", 1)[0]
        return cut + "\n\n_(reply trimmed at CHAT_MAX_REPLY_CHARS limit)_"
    return text


# ============================================================
# NO-IDEA FLOW / MUSE HELPERS (for dashboard + onboarding)
# Rich curated + lightweight generation for web "Step into the Lab".
# Mirrors Oracle's deeper /generate_problem (which uses full personas + Library + LLM).
# ============================================================

NO_IDEA_CATEGORIES = [
    "surprise", "foundations", "frontiers", "interdisciplinary", "applied",
    "theoretical", "historical", "everyday", "personal", "library_fork"
]

CURATED_STARTERS = {
    "surprise": [
        "What open question in your field would benefit most from multiple specialist perspectives and sandbox experiments?",
        "How can structured disagreement between researchers produce better hypotheses than consensus-seeking dialogue?",
    ],
    "foundations": [
        "What core assumption in a familiar domain should be tested before building further theory on top of it?",
        "Which definitions in a contested field are imprecise enough to block rigorous experimentation?",
    ],
    "frontiers": [
        "Where do current methods break down on problems that are high-stakes, novel, or poorly specified?",
        "What edge case would reveal whether a popular approach is robust or merely convenient?",
    ],
    "interdisciplinary": [
        "What connection between two fields is discussed informally but has never been modeled and tested?",
        "How would a tool from one discipline change conclusions in another if applied rigorously?",
    ],
    "applied": [
        "What practical system could be improved by a small, well-designed experiment rather than more opinion?",
        "Which real-world failure mode deserves a reproducible analysis before anyone proposes a fix?",
    ],
    "theoretical": [
        "What abstract question rewards clear formalization even if immediate data is limited?",
        "How can a thought experiment be turned into code that produces inspectable artifacts?",
    ],
    "historical": [
        "What past research program succeeded or failed for reasons we still misunderstand?",
        "Which old debate would change if revisited with modern tools and a multi-agent review loop?",
    ],
    "everyday": [
        "What seemingly simple phenomenon becomes interesting when you demand measurement and a testable model?",
        "Which everyday observation has a mechanism that specialists explain differently?",
    ],
    "personal": [
        "What question from lived experience could be reframed as a rigorous, testable research problem?",
        "How might family history or local context suggest a problem the squad can attack without losing rigor?",
    ],
    "library_fork": [
        "Building on a prior loop in this lab, what follow-up experiment would stress-test the last synthesis?",
        "Which conclusion from past artifacts should be revisited because new methods or data are available?",
    ],
}

def generate_no_idea_starter(category="surprise", family_note=""):
    """Lightweight generator for web UI and CLI fallback. Returns problem + rich metadata."""
    cat = category if category in NO_IDEA_CATEGORIES else "surprise"
    base = random.choice(CURATED_STARTERS.get(cat, CURATED_STARTERS["surprise"]))
    
    if family_note and len(family_note.strip()) > 3 and cat in ("family_history", "human_mysteries", "surprise"):
        base = base.replace("ancestral", f"ancestral (echoing {family_note[:70]})") if "ancestral" in base else \
               f"{base.rstrip('.')} — lightly resonant with themes of {family_note[:80]}."
    
    persona = random.choice(list(AGENT_PERSONAS.keys()))
    note = f"Curated starter from the No-Idea Portal ({cat}). Voiced through {persona}'s lens. The full Oracle Muse can deepen this."
    
    return {
        "problem_statement": base,
        "category": cat,
        "persona_hint": persona,
        "muse_note": note,
        "source": "lab_curated",
        "family_infused": bool(family_note and len(family_note.strip()) > 3)
    }

def get_recent_loops(limit=6):
    rows = _query_chimera_db(
        "SELECT loop_id, start_time, substr(problem_statement,1,140), CASE WHEN consensus_artefact IS NOT NULL THEN 1 ELSE 0 END FROM agora_loops ORDER BY loop_id DESC LIMIT ?",
        (limit,)
    )
    loops = []
    for lid, t, prob, has_synth in rows:
        # Enrich with artifact count
        art_count = 0
        for (aj,) in _query_lab_db("SELECT artifacts_json FROM lab_experiments WHERE loop_id = ?", (str(lid),)):
            try:
                art_count += len(json.loads(aj)) if aj else 0
            except: pass
        loops.append({
            "id": lid,
            "time": t,
            "problem": prob + "..." if prob and len(prob) > 130 else prob,
            "has_synthesis": bool(has_synth),
            "artifacts": art_count
        })
    # No cross-department placeholders — empty labs stay empty
    # GETAILABLIBRARY ENRICHMENT: page counts from the real codex if available
    if GETAILAB_LIBRARY_ENABLED and GETAILAB_LIB:
        for lp in loops:
            try:
                ps = GETAILAB_LIB.get_loop_as_pages(lp["id"])
                if ps:
                    lp["library_pages"] = len(ps)
                    lp["has_library_book"] = True
            except Exception:
                pass
    return loops

def _loop_report_paths(loop_id: int):
    """Search only this lab's report directories — never another department's loop_*.md."""
    paths = []
    try:
        reports_root = lab_reports_dir(ACTIVE_LAB_ID)
        paths.append(os.path.join(str(reports_root), f"loop_{loop_id}_report.md"))
        if is_chimera_lab(ACTIVE_LAB_ID):
            paths.append(os.path.join(ROOT_DIR, f"loop_{loop_id}_report.md"))
            paths.append(os.path.join(ROOT_DIR, "docs", "loops", f"loop_{loop_id}_report.md"))
    except Exception:
        paths.append(os.path.join(ROOT_DIR, f"loop_{loop_id}_report.md"))
    return paths


def _load_loop_report_md(loop_id: int) -> str:
    for path in _loop_report_paths(loop_id):
        if os.path.exists(path):
            with open(path, encoding="utf-8", errors="replace") as fh:
                return fh.read()
    return ""


def _parse_report_sections(report_md: str) -> Dict[str, Dict[str, str]]:
    """Extract per-scientist hypothesis/experiment blocks from a loop report markdown file."""
    import re
    sections: Dict[str, Dict[str, str]] = {}
    if not report_md:
        return sections
    for match in re.finditer(r"^##\s+(.+?)\s*$", report_md, re.M):
        heading = match.group(1).strip()
        start = match.end()
        next_match = re.search(r"^##\s+", report_md[start:], re.M)
        body = report_md[start: start + next_match.start()] if next_match else report_md[start:]
        body = body.strip()
        if "'s Hypothesis" in heading:
            name = heading.split("'s")[0].strip().lower()
            sections.setdefault(name, {})["hypothesis"] = body
        elif "'s Experiment" in heading:
            name = heading.split("'s")[0].strip().lower()
            sections.setdefault(name, {})["experiment_md"] = body
        elif heading == "Oracle's Consensus Artefact":
            sections["_oracle"] = {"synthesis": body}
    return sections


def _loop_scientist_outputs(loop_id: int) -> List[Dict[str, Any]]:
    """Full per-scientist code, stdout, stderr, and artifacts from lab_experiments."""
    outputs: List[Dict[str, Any]] = []
    rows = _query_lab_db(
        """
        SELECT agent_name, code, stdout, stderr, success, execution_time_ms, artifacts_json, experiment_name
        FROM lab_experiments
        WHERE loop_id = ?
        ORDER BY id ASC
        """,
        (str(loop_id),),
    )
    for agent_name, code, stdout, stderr, success, exec_ms, artifacts_json, experiment_name in rows:
        artifacts: List[str] = []
        try:
            artifacts = json.loads(artifacts_json or "[]") or []
        except Exception:
            pass
        outputs.append({
            "agent": (agent_name or "unknown").lower(),
            "experiment_name": experiment_name or "",
            "code": code or "",
            "stdout": stdout or "",
            "stderr": stderr or "",
            "success": bool(success),
            "execution_time_ms": exec_ms or 0,
            "artifacts": artifacts,
        })
    return outputs


def get_loop_detail(loop_id, full: bool = False):
    prob = ""
    synth = ""
    rows = _query_chimera_db(
        "SELECT problem_statement, consensus_artefact FROM agora_loops WHERE loop_id = ?",
        (loop_id,),
    )
    if rows:
        prob, synth = rows[0]

    art_count = 0
    for (aj,) in _query_lab_db("SELECT artifacts_json FROM lab_experiments WHERE loop_id = ?", (str(loop_id),)):
        try:
            art_count += len(json.loads(aj) or [])
        except Exception:
            pass

    report_md = _load_loop_report_md(loop_id)
    report_sections = _parse_report_sections(report_md)
    lab_outputs = _loop_scientist_outputs(loop_id)

    hypotheses: Dict[str, str] = {}
    if GETAILAB_LIBRARY_ENABLED and GETAILAB_LIB:
        try:
            for page in GETAILAB_LIB.get_loop_as_pages(loop_id):
                if page.page_type == "hypothesis" and page.agent:
                    hypotheses[page.agent.lower()] = page.content
        except Exception:
            pass
    for name, data in report_sections.items():
        if name.startswith("_"):
            continue
        if data.get("hypothesis") and name not in hypotheses:
            hypotheses[name] = data["hypothesis"]

    scientists: List[Dict[str, Any]] = []
    seen = set()
    for row in lab_outputs:
        agent = row["agent"]
        seen.add(agent)
        scientists.append({
            **row,
            "hypothesis": hypotheses.get(agent, ""),
        })
    for agent, hyp in sorted(hypotheses.items()):
        if agent not in seen:
            scientists.append({
                "agent": agent,
                "experiment_name": "",
                "code": "",
                "stdout": "",
                "stderr": "",
                "success": None,
                "execution_time_ms": 0,
                "artifacts": [],
                "hypothesis": hyp,
            })

    synthesis_full = synth or report_sections.get("_oracle", {}).get("synthesis", "")
    if not synthesis_full and report_md:
        synthesis_full = report_md

    if not prob and report_md:
        for line in report_md.splitlines():
            if line.lower().startswith("**problem:"):
                prob = line.split(":", 1)[-1].strip().strip("*")
                break

    if not prob and not report_md:
        return {
            "id": loop_id,
            "problem": "No record found for this loop.",
            "synthesis": (
                f"No data in {ACTIVE_LAB_ID} vault yet. "
                f"Artifacts: {ARTIFACTS_DIR} · Loop DB: {_agora_db_path()}"
            ),
            "synthesis_full": "",
            "artifacts": 0,
            "scientists": [],
            "has_full_report": False,
        }

    synthesis_display = synthesis_full if full else (synthesis_full[:1800] + ("…" if len(synthesis_full) > 1800 else ""))

    payload = {
        "id": loop_id,
        "problem": prob or "Research loop",
        "synthesis": synthesis_display,
        "synthesis_full": synthesis_full,
        "artifacts": art_count,
        "scientists": scientists,
        "scientist_count": len(scientists),
        "has_full_report": bool(report_md),
        "report_path": next((p for p in _loop_report_paths(loop_id) if os.path.exists(p)), None),
    }
    if full:
        payload["report_md"] = report_md
    return payload


# ============================================
# GETAILAB LIBRARY: User Notes / Directives Loader
# Ties the dashboard directly to your recorded problem statements,
# loop reports, and architectural philosophy (README + loop_*.md)
# Pure inspiration fuel. The field references *you*.
# ============================================

def load_user_directives(limit=6):
    """Load architect's notes, problem vectors, and key directives from the living record.
    These are YOUR words and vision — surfaced to pump tyres and ground every visit."""
    directives = []
    # Recent loop problems from DB are strongest (user-directed)
    try:
        rows = _query_chimera_db(
            "SELECT loop_id, start_time, problem_statement FROM agora_loops ORDER BY loop_id DESC LIMIT ?",
            (limit + 2,)
        )
        for lid, t, prob in rows:
            if prob and len(prob) > 20:
                excerpt = prob[:260].strip()
                if len(prob) > 260:
                    excerpt += "…"
                directives.append({
                    "loop": lid,
                    "time": t,
                    "note": excerpt,
                    "type": "directive",
                    "source": f"Loop #{lid} problem vector"
                })
    except Exception:
        pass

    # Supplement with loop reports from *this lab only* (forged labs never read the example lab root reports)
    report_files = []
    if is_chimera_lab(ACTIVE_LAB_ID):
        try:
            for f in os.listdir(ROOT_DIR):
                if f.startswith("loop_") and f.endswith("_report.md"):
                    report_files.append(os.path.join(ROOT_DIR, f))
            report_files.sort(key=lambda p: int(''.join(filter(str.isdigit, os.path.basename(p))) or 0), reverse=True)
        except Exception:
            report_files = []
    else:
        try:
            reports_root = str(lab_reports_dir(ACTIVE_LAB_ID))
            if os.path.isdir(reports_root):
                for f in os.listdir(reports_root):
                    if f.startswith("loop_") and f.endswith("_report.md"):
                        report_files.append(os.path.join(reports_root, f))
                report_files.sort(key=lambda p: int(''.join(filter(str.isdigit, os.path.basename(p))) or 0), reverse=True)
        except Exception:
            report_files = []

    for path in report_files[:3]:
        try:
            with open(path, 'r', encoding='utf-8', errors='ignore') as fh:
                content = fh.read()
            # Extract the Problem line + a juicy follow-up sentence or hypothesis hook
            lines = [l.strip() for l in content.split('\n') if l.strip()]
            prob_line = next((l for l in lines if l.lower().startswith("**problem:") or "problem:" in l.lower()), None)
            if not prob_line:
                prob_line = lines[2] if len(lines) > 2 else lines[0]
            # Grab a signature inspiring sentence from the architect's framing (often after the problem)
            insight = ""
            for ln in lines[3:12]:
                if len(ln) > 55 and any(k in ln.lower() for k in ["we are", "you are", "requires", "abandon", "building", "landscape", "vessel", "curvature", "manifold"]):
                    insight = ln[:210]
                    break
            directives.append({
                "loop": int(''.join(filter(str.isdigit, os.path.basename(path))) or 0),
                "time": "report",
                "note": (prob_line.replace("**Problem:**", "").replace("**problem:**", "").strip() + (" — " + insight if insight else ""))[:280],
                "type": "architect_note",
                "source": os.path.basename(path)
            })
        except Exception:
            continue

    try:
        if not is_chimera_lab(ACTIVE_LAB_ID):
            raise FileNotFoundError("skip readme for forged labs")
        readme_path = os.path.join(ROOT_DIR, "README.md")
        if os.path.exists(readme_path):
            with open(readme_path, 'r', encoding='utf-8', errors='ignore') as fh:
                r = fh.read()
            if "dialectic loop" in r.lower() or "research loops" in r.lower():
                directives.append({
                    "loop": 0,
                    "time": "reference",
                    "note": "Dialectic loop: hypothesis → implement → execute → synthesize → direction picker. Per-stage job tickets.",
                    "type": "reference",
                    "source": "README"
                })
    except Exception:
        pass

    # Dedup + limit + shuffle lightly for freshness
    seen = set()
    unique = []
    for d in directives:
        key = (d.get("loop"), d["note"][:40])
        if key not in seen:
            seen.add(key)
            unique.append(d)
    random.shuffle(unique)
    return unique[:limit]


def get_research_progress(stats, squad_size: int = 1):
    """Derived activity score from real loop and artifact counts. 0-100."""
    loops = stats.get("loops_completed", 0)
    arts = stats.get("total_artifacts", 0)
    books = stats.get("books_created", 0)
    squad = max(1, squad_size)
    balance = min(1.0, len(stats.get("agent_contributions", {})) / squad)
    score = min(98, int((loops * 4) + (arts / 20) + (books * 3) + (balance * 15)))
    return max(0, score)

def get_loop_full_for_export(loop_id):
    """Pull full research book/page content from DB (consensus) + MD reports for GetAiLabLibrary export.
    NOW INTEGRATED: Prefers structured pages from GetAiLabLibrary (books/pages from loops).
    Added via library_api prepare_loop_export_data (reports integration in library module).
    Legacy MD+DB fallback if library disabled."""
    # Prefer GetAiLabLibrary research books/pages for canonical export data
    if GETAILAB_LIBRARY_ENABLED and GETAILAB_LIB:
        try:
            lib_data = GETAILAB_LIB.prepare_loop_export_data(loop_id)
            if lib_data and lib_data.get("page_count", 0) > 0:
                # Enrich with legacy detail for artifacts if needed
                detail = get_loop_detail(loop_id)
                lib_data.setdefault("problem", detail.get("problem", lib_data.get("problem", "")))
                lib_data.setdefault("synthesis", detail.get("synthesis", lib_data.get("synthesis", "")))
                return lib_data
        except Exception as _lib_e:
            print(f"[GetAiLabLibrary export] prepare failed, fallback: {_lib_e}")

    # Legacy path (this lab's MD reports + DB only)
    detail = get_loop_detail(loop_id)
    full_report = detail.get('synthesis', '')
    for md_path in _loop_report_paths(loop_id):
        if os.path.exists(md_path):
            try:
                with open(md_path, 'r', encoding='utf-8', errors='replace') as f:
                    full_report = f.read()
                break
            except Exception:
                pass
    # enrich artifacts list
    artifacts_list = []
    for (aj,) in _query_lab_db("SELECT artifacts_json FROM lab_experiments WHERE loop_id = ?", (str(loop_id),)):
        try:
            if aj:
                artifacts_list.extend(json.loads(aj) or [])
        except: pass
    return {
        "id": loop_id,
        "problem": detail.get("problem", ""),
        "synthesis": detail.get("synthesis", ""),
        "full_report": full_report,
        "artifacts": detail.get("artifacts", 0),
        "artifacts_list": list(set(artifacts_list))[:30]  # dedup for slide sanity
    }

# --- PDF Generator (Entry Level / Free Tier) ---
def generate_pdf_export(loop_data):
    if FPDF is None:
        raise RuntimeError("fpdf2 not available. Install for PDF exports.")
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    pdf.set_text_color(30, 30, 40)
    # Sanitize em-dashes etc for Helvetica core font
    title = f"GetAiLab Research Book - Loop #{loop_data['id']}".replace("—", "-").replace("–", "-")
    pdf.cell(0, 12, title, ln=True, align='C')
    pdf.set_font("Helvetica", "", 10)
    pdf.set_text_color(100, 100, 120)
    pdf.cell(0, 8, "GetAiLab | GetAiLab Research Export", ln=True, align='C')
    pdf.ln(4)
    pdf.set_draw_color(252, 211, 77)
    pdf.line(10, pdf.get_y(), 200, pdf.get_y())
    pdf.ln(6)

    # Problem
    pdf.set_font("Helvetica", "B", 12)
    pdf.set_text_color(30, 30, 40)
    pdf.cell(0, 8, "PROBLEM VECTOR", ln=True)
    pdf.set_font("Helvetica", "", 10)
    pdf.set_text_color(50, 50, 60)
    pdf.multi_cell(0, 5, str(loop_data.get('problem', ''))[:1200].replace("—", "-"))
    pdf.ln(4)

    # Synthesis / Consensus
    pdf.set_font("Helvetica", "B", 12)
    pdf.set_text_color(30, 30, 40)
    pdf.cell(0, 8, "ORACLE CONSENSUS ARTEFACT", ln=True)
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(50, 50, 60)
    synth = str(loop_data.get('synthesis', loop_data.get('full_report', '')))[:3500].replace("—", "-")
    pdf.multi_cell(0, 4.5, synth)
    pdf.ln(4)

    # Artifacts summary
    pdf.set_font("Helvetica", "B", 12)
    pdf.set_text_color(30, 30, 40)
    pdf.cell(0, 8, f"ARTIFACTS VAULT ({loop_data.get('artifacts', 0)} files)", ln=True)
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(50, 50, 60)
    arts = ", ".join(loop_data.get('artifacts_list', [])[:12]) or "See vault for plots, CSVs, JSONs."
    pdf.multi_cell(0, 4.5, arts.replace("—", "-"))
    pdf.ln(3)

    pdf.set_font("Helvetica", "I", 8)
    pdf.set_text_color(120, 120, 140)
    pdf.cell(0, 6, "Generated by GetAiLab — loop export.", ln=True, align='C')

    out_path = os.path.join(ARTIFACTS_DIR, f"export_loop_{loop_data['id']}.pdf")
    pdf.output(out_path)
    return out_path

# --- PowerPoint Generator (optional, requires python-pptx) ---
def generate_pptx_export(loop_data):
    global Presentation, Inches, Pt, RgbColor, PP_ALIGN
    if Presentation is None:
        try:
            from pptx import Presentation as _P
            from pptx.util import Inches as _I, Pt as _Pt
            from pptx.dml.color import RGBColor as _Rgb
            from pptx.enum.text import PP_ALIGN as _PP
            Presentation, Inches, Pt, RgbColor, PP_ALIGN = _P, _I, _Pt, _Rgb, _PP
        except Exception as _imp_e:
            raise RuntimeError(f"python-pptx not available ({_imp_e}). Install for Pro PPTX exports.")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"GetAiLab Research Book: Loop #{loop_data['id']}"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RgbColor(252, 211, 77)
    p.alignment = PP_ALIGN.CENTER
    sub = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.8))
    sub.text_frame.paragraphs[0].text = "GetAiLab • GetAiLab Research Export"
    sub.text_frame.paragraphs[0].font.size = Pt(18)
    sub.text_frame.paragraphs[0].font.color.rgb = RgbColor(192, 132, 252)

    # Slide 2: Problem
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "PROBLEM VECTOR"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(103, 232, 249)
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = body.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(loop_data.get('problem', ''))
    p.font.size = Pt(16)
    p.font.color.rgb = RgbColor(224, 231, 255)

    # Slide 3: Consensus Artefact (core research book content)
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ORACLE CONSENSUS ARTEFACT"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(192, 132, 252)
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    synth = str(loop_data.get('synthesis', loop_data.get('full_report', '')))[:2800]
    for i, line in enumerate(synth.split('\n')[:35]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line[:140]
        p.font.size = Pt(13)
        p.font.color.rgb = RgbColor(224, 231, 255)

    # Slide 4: Key Artifacts & Vault
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"ARTIFACTS VAULT ({loop_data.get('artifacts', 0)} files)"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(103, 232, 249)
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = body.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Exported pages from the GetAiLab Library:"
    p.font.size = Pt(16)
    for art in loop_data.get('artifacts_list', [])[:18]:
        p = tf.add_paragraph()
        p.text = f"• {art}"
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(224, 231, 255)

    # Slide 5: Footer / Tier note
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Exported from GetAiLab loop records and artifact vault."
    p.font.size = Pt(22)
    p.font.color.rgb = RgbColor(252, 211, 77)
    p.alignment = PP_ALIGN.CENTER

    out_path = os.path.join(ARTIFACTS_DIR, f"export_loop_{loop_data['id']}.pptx")
    prs.save(out_path)
    return out_path

@app.route('/health', methods=['GET'])
def health_check():
    libs = ['numpy', 'scipy', 'matplotlib', 'pandas', 'sympy']
    try:
        import pyarrow  # noqa: F401
        libs.append('pyarrow')
    except ImportError:
        pass
    return jsonify({
        'agent': 'lab_sandbox_v2', 
        'status': 'active',
        'libraries': libs,
        'tools': ['execute', 'vision/extract', 'web/read', 'literature/search'],
        'workspace': ARTIFACTS_DIR
    })

def _sandbox_sanitize_text(text: str, limit: int = 80_000) -> str:
    """Keep report-safe stdout/stderr (strip NULs / lone surrogates from binary noise)."""
    if not text:
        return ""
    if isinstance(text, bytes):
        text = text.decode("utf-8", errors="replace")
    # NULs break markdown / some editors treat file as binary
    text = text.replace("\x00", "")
    text = text.encode("utf-8", errors="replace").decode("utf-8", errors="replace")
    if len(text) > limit:
        text = text[:limit] + "\n...[truncated by lab sandbox]...\n"
    return text


def _sandbox_rewrite_code(code: str) -> str:
    """
    Light rewrites for recurring LLM code smells that fail in the sandbox.
    Does NOT change science logic — only broken imports / path constants.
    """
    if not code:
        return code
    out = code
    # sympy: commutator lives under physics.quantum
    out = re.sub(
        r"from\s+sympy\s+import\s+\(([^)]*)\)",
        lambda m: _rewrite_sympy_import_paren(m.group(0), m.group(1)),
        out,
        flags=re.DOTALL,
    )
    out = re.sub(
        r"from\s+sympy\s+import\s+([^\n#]+)",
        lambda m: _rewrite_sympy_import_flat(m.group(0), m.group(1)),
        out,
    )
    out = out.replace(
        "from sympy import commutator",
        "from sympy.physics.quantum import Commutator as commutator",
    )
    out = out.replace(
        "from sympy import Commutator",
        "from sympy.physics.quantum import Commutator",
    )
    out = out.replace(
        "from sympy import commutation",
        "from sympy.physics.quantum import Commutator as commutation",
    )
    # common LLM typos — scipy has no commutor
    out = re.sub(
        r"from\s+scipy\.linalg\s+import\s+commutor\b",
        "def commutor(A, B):\n    import numpy as _np\n    A, B = _np.asarray(A), _np.asarray(B)\n    return A @ B - B @ A",
        out,
    )
    out = re.sub(
        r"from\s+scipy\s+import\s+commutor\b",
        "def commutor(A, B):\n    import numpy as _np\n    A, B = _np.asarray(A), _np.asarray(B)\n    return A @ B - B @ A",
        out,
    )
    out = re.sub(r"\bweightsnorm\b", "weights_norm", out)
    # hard-coded chimera / absolute workspaces → relative sandbox (loop workspace)
    out = re.sub(
        r"['\"]/(?:home/[^'\"]+/)?chimera/artifacts(?:/[^'\"]*)?['\"]",
        "'.'",
        out,
    )
    out = re.sub(
        r"['\"]/chimera/artifacts(?:/[^'\"]*)?['\"]",
        "'.'",
        out,
    )
    # /home/.../chimera_workspace (Loop 9 escape hatch)
    out = re.sub(
        r"['\"]/home/[^'\"]+/chimera_workspace(?:/[^'\"]*)?['\"]",
        "'.'",
        out,
    )
    out = re.sub(
        r"['\"][^'\"]*chimera_workspace[^'\"]*['\"]",
        "'.'",
        out,
    )
    out = re.sub(
        r"(OUTPUT_DIR|workspace|WORKSPACE|ARTIFACT_DIR|out_dir|OUT_DIR)\s*=\s*['\"]/[^'\"]+['\"]",
        r"\1 = '.'",
        out,
    )
    out = re.sub(
        r"(OUTPUT_DIR|workspace|WORKSPACE|ARTIFACT_DIR|out_dir|OUT_DIR)\s*=\s*['\"][^'\"]*chimera[^'\"]*['\"]",
        r"\1 = '.'",
        out,
    )
    return out


def _rewrite_sympy_import_flat(full: str, names: str) -> str:
    parts = [p.strip() for p in names.split(",") if p.strip()]
    keep, need_comm = [], False
    for p in parts:
        base = p.split(" as ")[0].strip()
        if base in ("commutator", "Commutator", "commutation"):
            need_comm = True
        else:
            keep.append(p)
    lines = []
    if keep:
        lines.append("from sympy import " + ", ".join(keep))
    if need_comm:
        lines.append("from sympy.physics.quantum import Commutator as commutator")
        lines.append("Commutator = commutator  # noqa: F841")
    return "\n".join(lines) if lines else full


def _rewrite_sympy_import_paren(full: str, names: str) -> str:
    return _rewrite_sympy_import_flat(full, names.replace("\n", " "))


_SANDBOX_PRELUDE = r'''# --- GetAiLab sandbox prelude (auto-injected; do not edit) ---
import os as _os, sys as _sys, json as _json, pathlib as _pathlib, builtins as _builtins
_WORK = _pathlib.Path(_os.environ.get("GETAILAB_ARTIFACT_DIR", ".")).resolve()
_os.chdir(_WORK)
_os.makedirs(_WORK, exist_ok=True)
for _sub in ("jungian_analysis", "artifacts", "output", "results", "data", "experiment_outputs"):
    _os.makedirs(_WORK / _sub, exist_ok=True)
# Alias legacy absolute paths → this loop workspace (so open/write stays in the vault)
for _legacy in (
    "/chimera/artifacts",
    "/home/deadly/x/chimera/artifacts",
    "/home/deadly/x/chimera_workspace",
    str(_pathlib.Path.home() / "chimera" / "artifacts"),
    str(_pathlib.Path.home() / "chimera_workspace"),
):
    try:
        parent = _os.path.dirname(_legacy) or "/"
        _os.makedirs(parent, exist_ok=True)
        if _os.path.islink(_legacy) or _os.path.exists(_legacy):
            continue
        _os.symlink(_WORK, _legacy, target_is_directory=True)
    except Exception:
        pass
# Auto-create parent dirs on write (fixes open("subdir/file.json","w") without makedirs)
_orig_open = _builtins.open
def _smart_open(file, mode="r", *a, **k):
    try:
        if isinstance(file, (str, _pathlib.Path)) and any(m in str(mode) for m in ("w", "a", "x")):
            parent = _pathlib.Path(file).parent
            if str(parent) not in ("", "."):
                parent.mkdir(parents=True, exist_ok=True)
    except Exception:
        pass
    return _orig_open(file, mode, *a, **k)
_builtins.open = _smart_open
# JSON safety for numpy / sympy / complex
_orig_dump, _orig_dumps = _json.dump, _json.dumps
def _json_default(o):
    try:
        import numpy as _np
        if isinstance(o, (_np.bool_,)):
            return bool(o)
        if isinstance(o, (_np.integer,)):
            return int(o)
        if isinstance(o, (_np.floating,)):
            return float(o)
        if isinstance(o, (_np.complexfloating, complex)):
            return {"re": float(o.real), "im": float(o.imag)}
        if isinstance(o, _np.ndarray):
            return o.tolist()
    except Exception:
        pass
    try:
        import sympy as _sp
        if isinstance(o, _sp.Basic):
            try:
                return float(o)
            except Exception:
                return str(o)
    except Exception:
        pass
    if isinstance(o, complex):
        return {"re": o.real, "im": o.imag}
    return str(o)
def _dump(obj, fp, *a, **k):
    k.setdefault("default", _json_default)
    return _orig_dump(obj, fp, *a, **k)
def _dumps(obj, *a, **k):
    k.setdefault("default", _json_default)
    return _orig_dumps(obj, *a, **k)
_json.dump, _json.dumps = _dump, _dumps
# Common science stack — fail loud only if venv is broken
import numpy, scipy, matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: F401
import pandas as pd  # noqa: F401
import sympy  # noqa: F401
try:
    from sympy.physics.quantum import Commutator as _Commutator
    if not hasattr(sympy, "commutator"):
        sympy.commutator = _Commutator
except Exception:
    pass
# scipy.linalg.commutor typo shim (LLMs invent this name)
try:
    import scipy.linalg as _sla
    if not hasattr(_sla, "commutor"):
        def _commutor(A, B):
            A, B = numpy.asarray(A), numpy.asarray(B)
            return A @ B - B @ A
        _sla.commutor = _commutor
except Exception:
    pass
# --- end prelude ---
'''


@app.route("/develop/scaffold", methods=["POST"])
def develop_scaffold():
    """Multi-file product scaffold (Engineering / Barry's crew DNA) — not single-file experiments."""
    data = request.get_json() or {}
    kind = (data.get("kind") or data.get("surface") or "web").strip()
    name = (data.get("name") or "project").strip()
    description = (data.get("description") or data.get("brief") or "").strip()
    agent = (data.get("agent_name") or data.get("agent") or "dev_shed").strip()
    do_zip = data.get("package", True)
    try:
        from getailab.develop import (
            create_build_manifest,
            create_project,
            list_project_files,
            package_project,
            write_files,
        )
        from getailab.develop.scaffolds import scaffold_for

        files = scaffold_for(kind, name, description=description)
        lab_id = os.environ.get("LAB_ID", "dev_shed")
        project = create_project(name, lab_id=lab_id)
        wr = write_files(files, project)
        create_build_manifest(
            project,
            agent=agent,
            build_type=kind,
            brief=description or name,
            extra={"kind": kind},
        )
        out = {
            "success": wr.get("success", False),
            "project_dir": str(project),
            "files_written": wr.get("files_written", []),
            "errors": wr.get("errors", []),
            "file_tree": list_project_files(project)[:100],
        }
        if do_zip and wr.get("success"):
            out["package"] = package_project(project, name)
        return jsonify(out), 200 if out["success"] else 400
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/develop/write_project", methods=["POST"])
def develop_write_project():
    """Write arbitrary multi-file dict into a new project dir (relative paths only)."""
    data = request.get_json() or {}
    name = (data.get("name") or "project").strip()
    files = data.get("files") or {}
    agent = (data.get("agent_name") or "dev_shed").strip()
    if not isinstance(files, dict) or not files:
        return jsonify({"success": False, "error": "files dict required"}), 400
    try:
        from getailab.develop import (
            create_build_manifest,
            create_project,
            list_project_files,
            package_project,
            write_files,
            write_to_product,
        )

        lab_id = os.environ.get("LAB_ID", "dev_shed")
        # Production line: land under product SoR when land_product=true
        if data.get("land_product") or data.get("to_product"):
            clean = {str(k): (v if isinstance(v, str) else str(v)) for k, v in files.items()}
            out = write_to_product(
                name,
                clean,
                source_lab=data.get("source_lab") or lab_id,
                agent=agent,
                overwrite=bool(data.get("overwrite", True)),
            )
            return jsonify(out), 200 if out.get("success") else 400

        project = create_project(name, lab_id=lab_id)
        # stringify values
        clean = {str(k): (v if isinstance(v, str) else str(v)) for k, v in files.items()}
        wr = write_files(clean, project)
        create_build_manifest(project, agent=agent, build_type="custom", brief=name)
        out = {
            "success": wr.get("success", False),
            "project_dir": str(project),
            "files_written": wr.get("files_written", []),
            "errors": wr.get("errors", []),
            "file_tree": list_project_files(project)[:100],
        }
        if data.get("package", True) and wr.get("success"):
            out["package"] = package_project(project, name)
        # optional promote to product SoR after build
        if data.get("promote_to_product") and wr.get("success"):
            from getailab.develop import promote_build_to_product

            out["product"] = promote_build_to_product(
                project, name, source_lab=data.get("source_lab") or lab_id, agent=agent
            )
        return jsonify(out), 200 if out["success"] else 400
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/develop/write_product", methods=["POST"])
def develop_write_product():
    """Write multi-file package straight into product SoR (Engineering)."""
    data = request.get_json() or {}
    name = (data.get("name") or data.get("package") or "package").strip()
    files = data.get("files") or {}
    agent = (data.get("agent_name") or data.get("agent") or "dev_shed").strip()
    if not isinstance(files, dict) or not files:
        return jsonify({"success": False, "error": "files dict required"}), 400
    try:
        from getailab.develop import write_to_product

        clean = {str(k): (v if isinstance(v, str) else str(v)) for k, v in files.items()}
        out = write_to_product(
            name,
            clean,
            source_lab=data.get("source_lab") or os.environ.get("LAB_ID"),
            agent=agent,
            overwrite=bool(data.get("overwrite", True)),
        )
        return jsonify(out), 200 if out.get("success") else 400
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/develop/review_product", methods=["POST", "GET"])
def develop_review_product():
    """Oracle / commander: test & review packages under product SoR."""
    data = request.get_json(silent=True) or {}
    if request.method == "GET":
        data = {**request.args.to_dict(), **data}
    package = (data.get("package") or data.get("name") or "").strip()
    run_smoke = str(data.get("run_smoke", "1")).lower() not in ("0", "false", "no")
    run_tests = str(data.get("run_tests", "1")).lower() not in ("0", "false", "no")
    try:
        from getailab.develop import (
            list_product_packages,
            review_all_products,
            review_product_package,
            product_root,
        )

        if package:
            out = review_product_package(
                package,
                source_lab=data.get("source_lab"),
                run_smoke=run_smoke,
                run_tests=run_tests,
            )
        elif data.get("list_only"):
            out = list_product_packages(source_lab=data.get("source_lab"))
        else:
            out = review_all_products(
                source_lab=data.get("source_lab"),
                run_smoke=run_smoke,
                run_tests=run_tests,
            )
        out["product_root"] = out.get("product_root") or str(product_root())
        return jsonify(out), 200
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route('/execute', methods=['POST'])
def execute():
    data = request.get_json() or {}
    code = data.get('code', '')
    agent_name = data.get('agent_name', 'unknown')
    loop_id = str(data.get('loop_id', 'global'))

    loop_workspace = os.path.join(ARTIFACTS_DIR, loop_id)
    os.makedirs(loop_workspace, exist_ok=True)

    if not code:
        return jsonify({'error': 'No code provided'}), 400

    rewritten = _sandbox_rewrite_code(code)
    # Inject prelude once (idempotent)
    if "GetAiLab sandbox prelude" not in rewritten:
        final_code = _SANDBOX_PRELUDE + "\n" + rewritten
    else:
        final_code = rewritten

    # Script stays at loop root so reports/synthesis still find exp_{agent}.py
    script_path = os.path.join(loop_workspace, f"exp_{agent_name}.py")
    with open(script_path, 'w', encoding='utf-8') as f:
        f.write(final_code)

    # Per-agent artifact isolation: prevent multi-agent stomping of results.csv etc.
    safe_agent = re.sub(r"[^\w.\-]+", "_", str(agent_name).strip()) or "unknown"
    agent_workspace = os.path.join(loop_workspace, safe_agent)
    os.makedirs(agent_workspace, exist_ok=True)

    env = os.environ.copy()
    env["GETAILAB_ARTIFACT_DIR"] = agent_workspace  # prelude chdirs here
    env["GETAILAB_LOOP_ID"] = str(loop_id)
    env["GETAILAB_AGENT"] = str(agent_name)
    env["GETAILAB_LOOP_WORKSPACE"] = loop_workspace
    env["MPLBACKEND"] = "Agg"
    # Prefer product packages on PYTHONPATH when vault has product/ next to artifacts/
    vault_root = os.path.dirname(ARTIFACTS_DIR)
    product_root = os.path.join(vault_root, "product")
    if os.path.isdir(product_root):
        env["GETAILAB_PRODUCT_ROOT"] = product_root
        env["PYTHONPATH"] = product_root + os.pathsep + env.get("PYTHONPATH", "")
    # Prefer the same interpreter that runs the lab (venv with scipy)
    python_exe = sys.executable

    start_time = time.time()
    proc_returncode = None
    try:
        result = subprocess.run(
            [python_exe, script_path],
            capture_output=True,
            text=True,
            timeout=1200,
            cwd=agent_workspace,
            env=env,
        )
        proc_returncode = result.returncode
        stdout = _sandbox_sanitize_text(result.stdout)
        stderr = _sandbox_sanitize_text(result.stderr)
        success, result_verdict = _evaluate_execute_success(result.returncode, stdout)
    except subprocess.TimeoutExpired:
        success, stdout, stderr = False, "", "CRITICAL: Execution timed out (1200s limit exceeded)."
        result_verdict = "timeout"
    except Exception as e:
        success, stdout, stderr = False, "", _sandbox_sanitize_text(str(e))
        result_verdict = "exception"

    exec_time = int((time.time() - start_time) * 1000)

    # List artifacts relative to loop_workspace so paths stay unique: {agent}/results.csv
    artifacts = []
    script_basename = os.path.basename(script_path)
    if os.path.isdir(agent_workspace):
        for root, _dirs, files in os.walk(agent_workspace):
            for name in files:
                full = os.path.join(root, name)
                rel = os.path.relpath(full, loop_workspace)
                artifacts.append(rel)
    # Optionally include loop-level files that belong to this agent (not exp scripts)
    for name in os.listdir(loop_workspace):
        full = os.path.join(loop_workspace, name)
        if not os.path.isfile(full):
            continue
        if name == script_basename or (name.startswith("exp_") and name.endswith(".py")):
            continue
        # Skip other agents' private dirs (already walked our own above)
        if name == safe_agent:
            continue
        # Only surface top-level non-agent files if they look agent-tagged
        if safe_agent in name or str(agent_name) in name:
            rel = os.path.relpath(full, loop_workspace)
            if rel not in artifacts:
                artifacts.append(rel)
    artifacts.sort()

    conn = sqlite3.connect(DB_PATH)
    conn.execute(
        "INSERT INTO lab_experiments (loop_id, agent_name, code, stdout, stderr, success, execution_time_ms, artifacts_json) "
        "VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
        (loop_id, agent_name, code, stdout, stderr, success, exec_time, json.dumps(artifacts)),
    )
    conn.commit()
    conn.close()

    _append_lab_ops({
        "event": "execute",
        "loop_id": loop_id,
        "agent": agent_name,
        "success": bool(success),
        "result_verdict": result_verdict,
        "execution_time_ms": exec_time,
        "returncode": proc_returncode,
        "stderr_head": (stderr or "")[:300],
        "stdout_tail": (stdout or "")[-200:],
        "artifact_count": len(artifacts),
    })

    return jsonify({
        'success': success,
        'stdout': stdout,
        'stderr': stderr,
        'execution_time_ms': exec_time,
        'artifacts': artifacts,
        'workspace_path': loop_workspace,
        'agent_workspace_path': agent_workspace,
        'result_verdict': result_verdict,
        'sandbox': {
            'python': python_exe,
            'prelude': True,
            'rewritten': rewritten != code,
        },
    })

@app.route('/api/llm/status')
def api_llm_status():
    """Show which LLM provider/model is active — useful for debugging configuration."""
    adapter = create_default_adapter()
    return jsonify({
        "config": get_env_provider_config(),
        "active": adapter.get_info(),
        "ready": adapter.is_ready(),
        "hint": "Set LLM_PROVIDER in .env (ollama, openai, google, anthropic, auto). Default is ollama.",
    })


@app.route('/vision/extract', methods=['POST'])
def vision_extract():
    """Sauron extraction endpoint used by scientists and run_lab.py."""
    data = request.get_json() or {}
    url = data.get('url', '')
    query = data.get('query', 'Extract key technical data')
    if not url:
        return jsonify({'success': False, 'error': 'url required'}), 400
    try:
        result = asyncio.run(SauronVision().extract(url, query))
        parsed = json.loads(result) if isinstance(result, str) else result
        ok = 'error' not in parsed
        return jsonify({'success': ok, 'data': parsed, 'text': result if ok else parsed.get('hint', '')})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/web/read', methods=['POST'])
def web_read():
    url = request.get_json().get('url', '')
    try:
        resp = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'}, timeout=30)
        soup = BeautifulSoup(resp.text, 'lxml')
        for tag in soup(['script', 'style']): tag.decompose()
        return jsonify({'success': True, 'text': md(str(soup), strip=['img'])[:20000]})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})


@app.route('/literature/search', methods=['POST'])
def literature_search():
    """Loop grounding — lab library/diary first; external APIs only if enabled."""
    data = request.get_json() or {}
    query = str(data.get('query') or data.get('problem_statement') or '').strip()
    if not query:
        return jsonify({'success': False, 'error': 'query or problem_statement required'}), 400
    sources = data.get('sources')
    max_per = int(data.get('max_per_source', 5))
    try:
        # Hot-reload so lab process picks up literature_search fixes without full reboot
        import importlib
        import getailab.literature_search as _lit
        importlib.reload(_lit)
        result = _lit.search_literature(query, sources=sources, max_per_source=max_per)
        ok = bool(result.get("success") or result.get("total", 0) > 0)
        return jsonify({
            'success': ok,
            'query': result.get('query'),
            'total': result.get('total', 0),
            'papers': result.get('papers', []),
            'results': result.get('results', {}),
            'formatted': result.get('formatted', ''),
            'errors': result.get('errors', []),
            'soft_warnings': result.get('soft_warnings', []),
            'sources_ok': result.get('sources_ok', []),
            'sources_attempted': result.get('sources_attempted', []),
            'mode': result.get('mode'),
            'prefer_library': result.get('prefer_library'),
            'allow_external': result.get('allow_external'),
            'real': True,
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e), 'real': True}), 500


# ============================================
# DASHBOARD ROUTES + LIVE PULSE ENGINE
# ============================================

# Classroom / university education API (curriculum, marking, web-started loops)
try:
    from getailab.classroom.web_api import bp as classroom_api_bp
    app.register_blueprint(classroom_api_bp)
except Exception as _classroom_api_err:
    print(f"  ⚠️  classroom API not loaded: {_classroom_api_err}")


@app.route('/')
@app.route('/dashboard')
@app.route('/lab')
def serve_dashboard():
    """Serve education-first dashboard (legacy research UI: /index.research.html)."""
    try:
        return send_from_directory(DASHBOARD_DIR, 'index.html')
    except Exception:
        # Graceful inline if file missing at boot (will be written by UI agent)
        return '<!doctype html><html><head><title>GetAiLab — Dashboard Loading</title></head><body style="background:#0a0a12;color:#e0e7ff;font-family:monospace;padding:40px"><h1 style="font-family:Georgia,serif;font-style:italic;color:#e8e4dc">e<sup>iπ</sup>+1=0</h1><p>GetAiLab Dashboard materializing. Refresh shortly.</p><p><a href="/api/stats" style="color:#fcd34d">Inspect live stats API</a></p></body></html>'


@app.route('/index.research.html')
def serve_research_dashboard():
    """Previous cosmic research UI (optional)."""
    try:
        return send_from_directory(DASHBOARD_DIR, 'index.research.html')
    except Exception:
        return jsonify({'error': 'Legacy UI not found'}), 404


@app.route('/manifest.json')
@app.route('/sw.js')
@app.route('/mobile_chat_stub.html')
def serve_dashboard_static():
    """PWA / mobile stubs next to the education dashboard."""
    name = request.path.lstrip('/')
    try:
        return send_from_directory(DASHBOARD_DIR, name)
    except Exception:
        return jsonify({'error': 'Not found'}), 404

@app.route('/api/config')
def api_config():
    """Frontend wiring: active lab, Oracle URL, squad — no hardcoded ports in dashboard."""
    try:
        from getailab.lab_config import get_lab_id, load_lab_config
        lid = get_lab_id()
        cfg = load_lab_config(lid)
    except Exception:
        lid = os.getenv("LAB_ID", "example")
        cfg = {}
    try:
        from personas.loader import get_squad_names
        squad = [n for n in get_squad_names() if n != "oracle"]
    except Exception:
        squad = list(AGENT_PERSONAS.keys())
    host = request.host.split(":")[0] if request.host else "localhost"
    lab_url = os.getenv("LAB_URL", f"http://{host}:{AGENT_PORT}").rstrip("/")
    try:
        paths = resolve_lab_paths(lid)
    except Exception:
        paths = dict(LAB_PATHS) if lid == ACTIVE_LAB_ID else {}
    return jsonify({
        "lab_id": lid,
        "display_name": cfg.get("display_name", lid.replace("_", " ").title()),
        "research_agenda": cfg.get("research_agenda", ""),
        "oracle_url": ORACLE_URL,
        "lab_url": lab_url,
        "lab_port": AGENT_PORT,
        "lab_results_db": DB_PATH,
        "artifacts_dir": ARTIFACTS_DIR,
        "vault_path": paths.get("vault"),
        "agora_db": paths.get("agora_db", _agora_db_path()),
        "reports_dir": paths.get("reports_dir"),
        "isolated": True,
        "squad": squad,
        "squad_size": len(squad),
        "agents": _refresh_agent_personas(),
        "library_enabled": bool(GETAILAB_LIBRARY_ENABLED and GETAILAB_LIB),
        "integrity_enabled": bool(GETAILAB_INTEGRITY_ENABLED),
        "tickets_enabled": bool(GETAILAB_TICKETS_ENABLED),
        "learning_enabled": bool(GETAILAB_LEARNING_ENABLED),
    })


@app.route('/api/stats')
def api_stats():
    stats = compute_real_stats()
    return jsonify(stats)

@app.route('/api/loops')
def api_loops():
    return jsonify(get_recent_loops())

@app.route('/api/reminders')
def api_reminders():
    stats = compute_real_stats()
    return jsonify(compute_reminders(stats))

@app.route('/api/inspire')
def api_inspire():
    agent = request.args.get("agent", "").strip().lower()
    return jsonify(get_random_inspiration(agent_filter=agent or None))


@app.route('/api/generate_starter')
def api_generate_starter():
    """Web-friendly No-Idea flow entry. Supports ?category=...&family_note=...
    Returns problem ready for preview or kickoff. Complements the deeper Oracle /generate_problem."""
    category = request.args.get('category', 'surprise')
    family_note = request.args.get('family_note', '') or request.args.get('family', '')
    starter = generate_no_idea_starter(category=category, family_note=family_note)
    return jsonify(starter)

@app.route('/api/ignite_muse', methods=['POST'])
def api_ignite_muse():
    """Onboarding portal: generate or accept a problem and return a CLI command to start a loop."""
    data = request.get_json(silent=True) or {}
    category = data.get('category', 'surprise')
    family_note = data.get('family_note', '') or data.get('family', '')
    problem = data.get('problem')

    if not problem:
        try:
            starter = generate_no_idea_starter(category=category, family_note=family_note)
            problem = starter.get('problem_statement', 'What research question should this squad tackle next?')
        except Exception:
            problem = "What experiment would best test a hypothesis generated by multiple specialist perspectives?"

    global _live_pulses
    pulse = {
        "ts": datetime.utcnow().isoformat(),
        "message": f"Starter problem selected: {problem[:85]}...",
        "boost": 1,
        "source": "muse_ignite"
    }
    _live_pulses.append(pulse)
    if len(_live_pulses) > 14:
        _live_pulses = _live_pulses[-14:]

    stats = compute_real_stats()
    stats["muse_ignition"] = {
        "problem": problem[:200],
        "category": category,
        "family_note": family_note[:60] if family_note else None,
        "note": "Run: python run_lab.py --problem \"...\" or POST to Oracle /initiate_loop"
    }

    cli_cmd = f'python3 run_lab.py --problem "{problem.replace(chr(34), chr(92)+chr(34))}"'
    return jsonify({
        "ok": True,
        "problem": problem,
        "category": category,
        "stats": stats,
        "pulse": pulse,
        "cli_command": cli_cmd,
        "next": "Start the full loop via CLI or Oracle /initiate_loop.",
    })

@app.route('/api/loop/<int:loop_id>')
def api_loop_detail(loop_id):
    full = request.args.get("full", "").lower() in ("1", "true", "yes")
    return jsonify(get_loop_detail(loop_id, full=full))

@app.route('/api/agents')
def api_agents():
    """Live squad from personas YAML — experiment counts from lab DB when available."""
    _refresh_agent_personas()
    exp_rows = _query_lab_db("SELECT agent_name, COUNT(*) FROM lab_experiments GROUP BY agent_name")
    exp_map = {r[0]: r[1] for r in exp_rows}
    return jsonify({
        k: {**v, "experiments": exp_map.get(k, exp_map.get(k.title(), 0))}
        for k, v in AGENT_PERSONAS.items()
    })

# NEW: GetAiLab Library Directives — your own recorded vision and notes.
# The dashboard now directly references and reveres the architect's problem statements.
@app.route('/api/directives')
def api_directives():
    dirs = load_user_directives(8)
    vault_sparks: List[Dict[str, Any]] = []
    if GETAILAB_LIBRARY_ENABLED and GETAILAB_LIB:
        try:
            from getailab.oracle.guardian import OracleGuardian
            root = Path(__file__).resolve().parent.parent / "data"
            guardian = OracleGuardian(GETAILAB_LIB.lab_id, root)
            topic = request.args.get("topic", "")
            vault_sparks = guardian.get_vault_sparks(topic=topic, limit=4)
        except Exception:
            pass
    merged = vault_sparks + dirs
    return jsonify({
        "directives": merged,
        "vault_sparks": vault_sparks,
        "count": len(merged),
    })


@app.route('/api/inspiration/vault')
def api_inspiration_vault():
    """Oracle-curated snippets from the library vault (user-safe)."""
    topic = request.args.get("q", "") or request.args.get("topic", "")
    if not GETAILAB_LIBRARY_ENABLED or not GETAILAB_LIB:
        return jsonify({"enabled": False, "sparks": []})
    try:
        from getailab.oracle.guardian import OracleGuardian
        root = Path(__file__).resolve().parent.parent / "data"
        guardian = OracleGuardian(GETAILAB_LIB.lab_id, root)
        return jsonify({
            "enabled": True,
            "topic": topic,
            "sparks": guardian.get_safe_inspiration_for_user(topic, {}),
            "cards": guardian.get_vault_sparks(topic=topic, limit=6),
        })
    except Exception as e:
        return jsonify({"enabled": False, "error": str(e), "sparks": []})

# User can log a reflection/note that joins the activity feed.
# This makes the UI truly two-way and personal. Reflections get pulsed.
@app.route('/api/reflect', methods=['POST'])
def log_reflection():
    data = request.get_json() or {}
    reflection = data.get('reflection', '').strip()
    if not reflection or len(reflection) < 3:
        return jsonify({"ok": False, "error": "Reflection too faint."}), 400
    global _live_pulses
    msg = f"Architect's reflection received: {reflection[:110]}"
    pulse = {
        "ts": datetime.utcnow().isoformat(),
        "message": msg,
        "boost": random.randint(2, 5),
        "source": "reflection"
    }
    _live_pulses.append(pulse)
    if len(_live_pulses) > 14:
        _live_pulses = _live_pulses[-14:]
    # Boost real stats a bit for the session
    from getailab.resonance import record_engagement
    record_engagement(ACTIVE_LAB_ID, boost=3, activity=True)
    stats = compute_real_stats()
    stats["pulse"] = pulse
    return jsonify({"ok": True, "pulse": pulse, "stats": stats, "message": "Reflection recorded."})

# Live field — SSE for those who can listen, plus simple pollable pulse
_live_pulses = []

@app.route('/api/pulse', methods=['POST'])
def trigger_pulse():
    """User or system can broadcast a field pulse. Returns fresh data + message."""
    global _live_pulses
    data = request.get_json() or {}
    msg = data.get('message', random.choice([
        "New experiment queued.",
        "Loop phase updated.",
        "Squad activity detected.",
        "Artifact vault updated."
    ]))
    # Occasionally tie pulse to a directive for extra personal delight
    if random.random() < 0.28:
        try:
            d = load_user_directives(1)[0]
            msg = f"Field resonates with your Loop #{d.get('loop','?')} note: {d['note'][:85]}..."
        except: pass
    pulse = {
        "ts": datetime.utcnow().isoformat(),
        "message": msg,
        "boost": random.randint(1, 3)
    }
    _live_pulses.append(pulse)
    if len(_live_pulses) > 12:
        _live_pulses = _live_pulses[-12:]
    from getailab.resonance import record_engagement
    record_engagement(ACTIVE_LAB_ID, boost=int(pulse.get("boost", 1)), activity=True)
    stats = compute_real_stats()
    stats["pulse"] = pulse
    return jsonify({"ok": True, "pulse": pulse, "stats": stats})


# NEW: Delightful interactive endpoint — "resonate with my note" from UI directives.
# Pumps tyres, logs to live field, returns fresh stats + affirmation. Ties user notes directly to action.
@app.route('/api/resonate', methods=['POST'])
def resonate_with_note():
    data = request.get_json() or {}
    note_ref = data.get('note_ref', 'your vision')
    loop_ref = data.get('loop', '?')
    boost = int(data.get('boost', 4))
    global _live_pulses
    msg = f"Architect resonated with L{loop_ref} note: {note_ref[:75]}... The field answers."
    pulse = {"ts": datetime.utcnow().isoformat(), "message": msg, "boost": boost, "source": "resonate_directive"}
    _live_pulses.append(pulse)
    if len(_live_pulses) > 14:
        _live_pulses = _live_pulses[-14:]
    from getailab.resonance import record_engagement
    record_engagement(ACTIVE_LAB_ID, boost=boost, activity=True)
    stats = compute_real_stats()
    return jsonify({"ok": True, "pulse": pulse, "stats": stats, "affirmation": "Note recorded in the activity feed."})


@app.route('/api/resonance/nudge', methods=['POST'])
def api_resonance_nudge():
    """Persist dashboard nudges (inspiration ring, streak clicks) — per-lab only."""
    data = request.get_json() or {}
    boost = int(data.get("boost") or data.get("amount") or 1)
    from getailab.resonance import record_engagement
    record_engagement(ACTIVE_LAB_ID, boost=max(0, min(12, boost)), activity=True)
    stats = compute_real_stats()
    return jsonify({"ok": True, "stats": stats})

@app.route('/api/live')
def api_live():
    """Simple pollable live events (and SSE if client uses EventSource)."""
    # Return recent pulses + a synthetic live heartbeat
    global _live_pulses
    if not _live_pulses:
        _live_pulses = [{"ts": datetime.utcnow().isoformat(), "message": "Activity monitor online.", "boost": 0}]
    heartbeat = {
        "ts": datetime.utcnow().isoformat(),
        "message": random.choice(["Lab idle.", "Checking loop records.", "Artifact vault reachable.", "Services nominal."]),
        "boost": 0
    }
    return jsonify({
        "pulses": _live_pulses[-6:],
        "heartbeat": heartbeat,
        "timestamp": datetime.utcnow().isoformat(),
        "live_directive": load_user_directives(1)[0] if random.random() > 0.5 else None
    })

@app.route('/api/stream')
def stream_live():
    """Server Sent Events endpoint for true liveness (clients: new EventSource('/api/stream'))."""
    def event_generator():
        import time as _time
        yield "data: {\"type\":\"connect\",\"msg\":\"Activity stream connected.\"}\n\n"
        for i in range(12):  # richer stream — more interactive sauce + Muse onboarding echoes
            _time.sleep(1.6)
            # Mix of generic + occasional DIRECT personal reference to user's own notes + real library
            if random.random() < 0.42:
                try:
                    d = load_user_directives(1)[0]
                    pulse_msg = f"Recent loop #{d.get('loop', '?')} noted in activity feed."
                except:
                    pulse_msg = "Library index updated."
            elif random.random() < 0.18:
                # NEW SAUCE: Occasionally stream a fresh Muse-generated starter to inspire "no idea" users live
                try:
                    cat = random.choice(NO_IDEA_CATEGORIES)
                    starter = generate_no_idea_starter(category=cat)
                    pulse_msg = f"Muse whispers a new resonance ({cat}): {starter['problem_statement'][:110]}... — Step into the portal."
                except:
                    pulse_msg = "A new curiosity vector is forming in the Library."
            else:
                pulse_msg = random.choice([
                    "New artifact saved.",
                    "Loop phase updated.",
                    "Squad contribution logged.",
                    "Library page indexed.",
                    "Research progress updated.",
                    "Experiment record written."
                ])
            payload = {
                "type": "pulse",
                "msg": pulse_msg,
                "t": datetime.utcnow().isoformat(),
                "library_tie": random.random() < 0.35,  # hints frontend to refresh directives
                "muse_echo": random.random() < 0.15
            }
            yield f"data: {json.dumps(payload)}\n\n"
            # Occasional real stats heartbeat for live dashboard updates
            if i % 3 == 0:
                try:
                    s = compute_real_stats()
                    hb = {"type":"heartbeat","stats_delta": {"inspiration": s.get("inspiration_score"), "library_pages": s.get("library_pages")}, "t": datetime.utcnow().isoformat()}
                    yield f"data: {json.dumps(hb)}\n\n"
                except: pass
    return app.response_class(event_generator(), mimetype='text/event-stream')

@app.route('/api/export/loop/<int:loop_id>')
def api_export_loop(loop_id):
    """Export loop summary as PDF or PPTX (if python-pptx installed)."""
    fmt = (request.args.get('format') or 'pdf').lower().strip()
    if fmt not in ('pdf', 'pptx'):
        fmt = 'pdf'

    try:
        loop_data = get_loop_full_for_export(loop_id)
        if fmt == 'pdf':
            path = generate_pdf_export(loop_data)
            fname = f"GetAiLab_Loop{loop_id}_ResearchBook.pdf"
            return send_from_directory(os.path.dirname(path), os.path.basename(path),
                                       as_attachment=True, download_name=fname)
        path = generate_pptx_export(loop_data)
        fname = f"GetAiLab_Loop{loop_id}_ResearchBook.pptx"
        return send_from_directory(os.path.dirname(path), os.path.basename(path),
                                   as_attachment=True, download_name=fname)
    except Exception as e:
        return jsonify({"error": f"Export failed: {str(e)}"}), 500

# ============================================================
# GETAILAB UNIFIED CHAT + MOBILE API (Cross-Platform: Web, CLI, Android, iOS)
# Real interactive council responses. Used by enhanced CLI, PWA chat, mobile stubs.
# ============================================================

@app.route('/api/chat', methods=['POST'])
def api_chat():
    """Primary chat endpoint for GetAiLab. Supports CLI, web dashboard, Android/iOS WebView or PWA.
    Returns council/agent voice response. Updates field on every interaction."""
    data = request.get_json() or {}
    user_msg = data.get('message', '').strip()
    source = data.get('source', 'web')
    plat = data.get('platform', 'unknown')
    history = data.get('history', [])[-3:]

    if not user_msg:
        return jsonify({'reply': 'Send a message to start the conversation.', 'agent': 'ORACLE'}), 200

    # Record pulse for live monitor (visible on web + mobile)
    global _live_pulses
    pulse_msg = f"Chat from {source}@{plat}: {user_msg[:60]}..."
    _live_pulses.append({"ts": datetime.utcnow().isoformat(), "message": pulse_msg, "boost": 2})
    if len(_live_pulses) > 12: _live_pulses = _live_pulses[-12:]

    agent_key = _pick_chat_agent(user_msg)
    agent = agent_key.upper() if agent_key != "oracle" else "ORACLE"
    reply = None

    # Loop-direction questions can still use Oracle's recommend_next
    if re.search(r"\b(next (loop|direction|research)|what should we (study|research|loop))\b", user_msg, re.I):
        try:
            rec = requests.post(ORACLE_URL + "/recommend_next", json={
                "synthesis": user_msg,
                "user_comment": user_msg,
            }, timeout=60).json()
            reply = (rec.get("oracle_rationale") or rec.get("recommendation") or "").strip()
            max_chars = int(os.getenv("CHAT_MAX_REPLY_CHARS", "8000"))
            if len(reply) > max_chars:
                reply = reply[:max_chars].rsplit(" ", 1)[0] + "\n\n_(reply trimmed)_"
            agent = "ORACLE"
        except Exception:
            pass

    if not reply:
        try:
            reply = _generate_council_reply(user_msg, history, agent_key)
        except Exception as exc:
            print(f"[CHAT] LLM fallback after error: {exc}")
            insp = get_random_inspiration()
            agent = insp.get("agent", "ORACLE").upper()
            reply = (
                f"(Council chat temporarily offline — {exc}) "
                f"{insp['quote']}"
            )

    # Nudge inspiration + live for all clients (mobile included)
    try:
        nudge_boost = 2 if source in ['cli', 'mobile', 'android', 'ios'] else 1
    except: pass

    learner_ctx = None
    if GETAILAB_LEARNING_ENABLED and get_adaptive_learner:
        try:
            user_id = data.get("user_id") or os.getenv("GETAILAB_USER_ID", "default")
            learner = get_adaptive_learner(user_id)
            learner.record_interaction(
                correct=None,
                response_quality=learner.estimate_message_quality(user_msg),
                subject=learner.extract_topic_hint(user_msg),
                concept=learner.extract_topic_hint(user_msg),
            )
            learner_ctx = learner.get_adaptive_context()
            intervention = learner_ctx.get("intervention")
            if intervention == "supportive" and reply:
                reply = f"{learner_ctx.get('coaching_hint', '')} {reply}"
        except Exception:
            pass

    payload = {
        'reply': reply,
        'agent': agent.upper(),
        'color': AGENT_PERSONAS.get(agent.lower(), {}).get('color', '#c084fc'),
        'timestamp': datetime.utcnow().isoformat(),
        'platform_echo': plat,
        'source': source,
        'field_pulse': True,
    }
    if learner_ctx:
        payload['learner'] = learner_ctx
    return jsonify(payload)

@app.route('/api/mobile/chat', methods=['POST', 'GET'])
def api_mobile_chat():
    """Dedicated lightweight chat stub endpoint optimized for Android/iOS WebView, React Native, or PWA.
    Returns minimal payload + suggestions for touch interfaces. Cross-platform parity with /api/chat."""
    if request.method == 'GET':
        return jsonify({
            'status': 'ready',
            'title': 'GetAiLab Council Mobile Chat',
            'agents': list(AGENT_PERSONAS.keys()),
            'instructions': 'POST message. Works in native WebViews on iOS/Android. Same backend as web+CLI.'
        })
    data = request.get_json() or {}
    data['source'] = data.get('source', 'mobile')
    data['platform'] = data.get('platform', 'mobile')
    chat_resp = api_chat()
    try:
        resp_json = chat_resp.get_json()
    except:
        resp_json = {'reply': 'Field resonance received.', 'agent': 'COUNCIL'}
    resp_json['quick_actions'] = [
        "Pulse the field",
        "Seed hypothesis from council",
        "Open library vault",
        "Start new dialectic loop"
    ]
    resp_json['mobile_optimized'] = True
    return jsonify(resp_json)

@app.route('/api/mobile/status')
def api_mobile_status():
    """Lightweight status for mobile apps / PWAs on Android or iOS."""
    stats = compute_real_stats()
    return jsonify({
        "getailab": "LIVE",
        "version": "V4",
        "lab_id": os.getenv("LAB_ID", "example"),
        "oracle_url": ORACLE_URL,
        "platforms": ["web", "windows", "macos", "linux", "android", "ios"],
        "full_support": "web + Win + macOS + Linux + Android + iOS (PWA + WebView stubs). CLI parity via --chat. Docker universal.",
        "stats": {
            "loops": stats.get("loops_completed"),
            "artifacts": stats.get("total_artifacts"),
            "resonance": stats.get("inspiration_score")
        },
        "chat_endpoint": "/api/mobile/chat",
        "chat_unified": "/api/chat",
        "dashboard": "/",
        "pwa_installable": True,
        "stubs": "mobile_chat_stub.html + JS bridges for native",
        "how_to_run_note": "Use host python run_lab.py --support for platform-specific commands. All clients identical."
    })

# End of GetAiLab chat + mobile cross-platform additions


# ============================================================
# GETAILABLIBRARY PUBLIC API ENDPOINTS (exposed here for dashboard + clients)
# Create / search books. View pages (hypotheses, code, artifacts, synthesis + checksums + provenance).
# Auto-integrated: every loop completion via orchestrator adds to library.
# Current outputs (artifacts, loops, synthesis) are first-class pages.
# Search is live & inspiring (serendipity + persona quotes injected).
# ============================================================

@app.route('/api/library/status')
def api_library_status():
    if not GETAILAB_LIBRARY_ENABLED or not GETAILAB_LIB:
        return jsonify({"enabled": False, "message": "Library integration pending — use /api/loops and /api/loop/<id> for now."})
    book = GETAILAB_LIB.get_or_create_default_book()
    summary = GETAILAB_LIB.get_recent_library_summary(6)
    return jsonify({
        "enabled": True,
        "book": {"id": book.book_id, "title": book.title, "slug": book.slug, "pages": len(book.page_ids)},
        "summary": summary,
        "vault_path": str(GETAILAB_LIB.persistence.root)
    })

@app.route('/api/library/books')
def api_library_books():
    if not GETAILAB_LIBRARY_ENABLED or not GETAILAB_LIB:
        return jsonify([])
    books = GETAILAB_LIB.list_books()
    return jsonify([{
        "book_id": b.book_id,
        "title": b.title,
        "slug": b.slug,
        "page_count": len(b.page_ids),
        "created": b.created_at,
        "inspiration": b.metadata.get("inspiration_score", 80)
    } for b in books])

@app.route('/api/library/search')
def api_library_search():
    q = request.args.get('q', '')
    limit = int(request.args.get('limit', 12))
    loop_filter = request.args.get('loop_id')
    filters = {"loop_id": int(loop_filter)} if loop_filter and loop_filter.isdigit() else None
    if not GETAILAB_LIBRARY_ENABLED or not GETAILAB_LIB:
        return jsonify({"results": [], "message": "Library offline — fall back to /api/loops and /api/loop/<id>"})
    hits = GETAILAB_LIB.search(q, filters=filters, limit=limit)
    return jsonify({"query": q, "results": hits, "count": len(hits)})

@app.route('/api/library/page/<page_id>')
def api_library_page(page_id):
    if not GETAILAB_LIBRARY_ENABLED or not GETAILAB_LIB:
        return jsonify({"error": "library disabled"}), 503
    page = GETAILAB_LIB.get_page(page_id)
    if not page:
        return jsonify({"error": "page not found"}), 404
    d = page.to_dict()
    d["inspiration"] = {"quote": "Registered page with checksum and provenance."}
    return jsonify(d)

@app.route('/api/library/loop/<int:loop_id>/pages')
def api_library_loop_pages(loop_id):
    """View the entire research loop reified as GetAiLab book pages with metadata, checksums, doccontrol tickets, artifacts."""
    if not GETAILAB_LIBRARY_ENABLED or not GETAILAB_LIB:
        return api_loop_detail(loop_id)
    pages = GETAILAB_LIB.get_loop_as_pages(loop_id)
    return jsonify({
        "loop_id": loop_id,
        "as_book_pages": len(pages),
        "pages": [p.to_dict(include_content=False) for p in pages]
    })

@app.route('/api/library/loop/<int:loop_id>/full')
def api_library_loop_full(loop_id):
    if not GETAILAB_LIBRARY_ENABLED or not GETAILAB_LIB:
        return api_loop_detail(loop_id)
    pages = GETAILAB_LIB.get_loop_as_pages(loop_id)
    return jsonify({
        "loop_id": loop_id,
        "page_count": len(pages),
        "pages": [p.to_dict() for p in pages]  # full: content + checksum + provenance + artifact refs
    })

@app.route('/api/library/verify')
def api_library_verify():
    if not GETAILAB_LIBRARY_ENABLED or not GETAILAB_LIB:
        return jsonify({"ok": False, "message": "library not active"})
    book_filter = request.args.get("book")
    include_indexes = request.args.get("indexes", "1").lower() not in ("0", "false", "no")
    if request.args.get("reindex", "").lower() in ("1", "true", "yes"):
        reindex_library(GETAILAB_LIB.lab_id)
    book = None
    if book_filter:
        from getailab.library.models import LibraryBook
        book = LibraryBook(
            book_id=book_filter.lower().strip(),
            title=book_filter.title(),
            slug=book_filter.lower().strip(),
        )
    else:
        book = GETAILAB_LIB.get_or_create_default_book()
    report = GETAILAB_LIB.verify_library_integrity(book, include_indexes=include_indexes)
    report["enabled"] = True
    return jsonify(report)


# ============================================================
# ADAPTIVE LEARNER API (Gabby user layer — education division pull)
# ============================================================

def _learner_user_id() -> str:
    return (
        request.args.get("user_id")
        or (request.get_json(silent=True) or {}).get("user_id")
        or os.getenv("GETAILAB_USER_ID", "default")
    )


@app.route('/api/learner/profile')
def api_learner_profile():
    if not GETAILAB_LEARNING_ENABLED or not get_adaptive_learner:
        return jsonify({"enabled": False, "message": "adaptive learner not active"}), 503
    user_id = _learner_user_id()
    learner = get_adaptive_learner(user_id)
    return jsonify({"enabled": True, **learner.get_adaptive_context()})


@app.route('/api/learner/interaction', methods=['POST'])
def api_learner_interaction():
    if not GETAILAB_LEARNING_ENABLED or not get_adaptive_learner:
        return jsonify({"enabled": False, "message": "adaptive learner not active"}), 503
    data = request.get_json() or {}
    user_id = data.get("user_id") or _learner_user_id()
    learner = get_adaptive_learner(user_id)
    message = data.get("message", "")
    quality = data.get("response_quality")
    if quality is None and message:
        quality = learner.estimate_message_quality(message)
    learner.record_interaction(
        correct=data.get("correct"),
        response_quality=float(quality or 0.5),
        subject=data.get("subject") or learner.extract_topic_hint(message),
        concept=data.get("concept"),
        mastered=bool(data.get("mastered")),
        struggling=bool(data.get("struggling")),
    )
    return jsonify({"enabled": True, **learner.get_adaptive_context()})


@app.route('/api/learner/coaching')
def api_learner_coaching():
    if not GETAILAB_LEARNING_ENABLED or not get_adaptive_learner:
        return jsonify({"enabled": False, "message": "adaptive learner not active"}), 503
    learner = get_adaptive_learner(_learner_user_id())
    return jsonify({
        "enabled": True,
        "message": learner.get_coaching_message(),
        **learner.get_adaptive_context(),
    })


@app.route('/api/learner/loop/<int:loop_id>', methods=['POST'])
def api_learner_loop_event(loop_id):
    """Record that the user observed or completed a loop."""
    if not GETAILAB_LEARNING_ENABLED or not get_adaptive_learner:
        return jsonify({"enabled": False, "message": "adaptive learner not active"}), 503
    data = request.get_json() or {}
    user_id = data.get("user_id") or _learner_user_id()
    learner = get_adaptive_learner(user_id)
    topic = data.get("topic", "")
    if data.get("completed"):
        learner.record_loop_completed(loop_id, topic=topic)
    else:
        learner.record_loop_observed(loop_id, topic=topic)
    return jsonify({"enabled": True, "loop_id": loop_id, **learner.get_adaptive_context()})


# ============================================================
# INTEGRITY API (Merkle snapshots + crush test — old-mate-og pull)
# ============================================================

def _integrity_lab_id() -> str:
    if GETAILAB_LIB:
        return GETAILAB_LIB.lab_id
    return request.args.get("lab_id") or os.getenv("LAB_ID", "example")


def _integrity_book_filter() -> str | None:
    book = request.args.get("book")
    if not book:
        return None
    book = book.lower().strip()
    return None if book in ("", "codex", "all") else book


@app.route('/api/integrity/status')
def api_integrity_status():
    if not GETAILAB_INTEGRITY_ENABLED or not merkle_status:
        return jsonify({"enabled": False, "message": "integrity API not active"})
    lab_id = _integrity_lab_id()
    status = merkle_status(lab_id)
    status["enabled"] = True
    return jsonify(status)


@app.route('/api/integrity/verify')
def api_integrity_verify():
    """Crush-test vault pages (+ optional SQLite index check)."""
    if not GETAILAB_INTEGRITY_ENABLED or not crush_test_vault:
        return jsonify({"enabled": False, "message": "integrity API not active"}), 503
    lab_id = _integrity_lab_id()
    book_id = _integrity_book_filter()
    include_indexes = request.args.get("indexes", "1").lower() not in ("0", "false", "no")
    if request.args.get("reindex", "").lower() in ("1", "true", "yes") and reindex_library:
        reindex_library(lab_id)

    pages = crush_test_vault(lab_id, book_id=book_id)
    report: Dict[str, Any] = {
        "enabled": True,
        "integrity": pages.get("integrity", "FAIL"),
        "lab_id": lab_id,
        "book_filter": book_id,
        "pages": pages,
    }
    if include_indexes and book_id is None and crush_test_indexes:
        indexes = crush_test_indexes(lab_id)
        report["indexes"] = indexes
        if indexes.get("integrity") != "PASS":
            report["integrity"] = "FAIL"
    report["ok"] = report["integrity"] == "PASS"
    return jsonify(report)


@app.route('/api/integrity/verify/full')
def api_integrity_verify_full():
    """Full report: crush test + indexes + Merkle status (optional live scan)."""
    if not GETAILAB_INTEGRITY_ENABLED or not verify_full:
        return jsonify({"enabled": False, "message": "integrity API not active"}), 503
    lab_id = _integrity_lab_id()
    book_id = _integrity_book_filter()
    merkle_scan = request.args.get("scan", "").lower() in ("1", "true", "yes")
    loop_id_raw = request.args.get("loop")
    loop_id = int(loop_id_raw) if loop_id_raw and loop_id_raw.isdigit() else None
    if request.args.get("reindex", "").lower() in ("1", "true", "yes") and reindex_library:
        reindex_library(lab_id)

    report = verify_full(
        lab_id,
        book_id=book_id,
        merkle_scan=merkle_scan,
        loop_id=loop_id,
    )
    report["enabled"] = True
    report["ok"] = report.get("integrity") == "PASS"
    return jsonify(report)


@app.route('/api/integrity/scan', methods=['POST'])
def api_integrity_scan():
    """
    Build Merkle snapshots for vault and/or lab artifacts.

    JSON body or query params: vault, artifacts, loop_id, lab_id, rotate (default true).
    """
    if not GETAILAB_INTEGRITY_ENABLED or not scan_integrity_targets:
        return jsonify({"enabled": False, "message": "integrity API not active"}), 503

    data = request.get_json(silent=True) or {}
    lab_id = data.get("lab_id") or _integrity_lab_id()
    loop_id = data.get("loop_id")
    if loop_id is None and request.args.get("loop", "").isdigit():
        loop_id = int(request.args.get("loop"))

    def _bool(key: str, default: bool) -> bool:
        if key in data:
            return bool(data[key])
        raw = request.args.get(key)
        if raw is None:
            return default
        return raw.lower() not in ("0", "false", "no")

    vault = _bool("vault", loop_id is None)
    artifacts = _bool("artifacts", loop_id is None)
    rotate_previous = _bool("rotate", True)

    if loop_id is not None:
        vault = _bool("vault", False)
        artifacts = _bool("artifacts", False)

    result = scan_integrity_targets(
        lab_id,
        vault=vault,
        artifacts=artifacts,
        loop_id=loop_id,
        rotate_previous=rotate_previous,
    )
    result["enabled"] = True
    return jsonify(result)


@app.route('/api/integrity/reindex', methods=['POST'])
def api_integrity_reindex():
    """Rebuild all scientist book + codex SQLite indexes."""
    if not GETAILAB_LIBRARY_ENABLED or not GETAILAB_LIB or not reindex_library:
        return jsonify({"enabled": False, "message": "library not active"}), 503
    lab_id = _integrity_lab_id()
    result = reindex_library(lab_id)
    result["enabled"] = True
    result["ok"] = True
    return jsonify(result)


@app.route('/api/integrity/sign/status')
def api_integrity_sign_status():
    if not GETAILAB_INTEGRITY_ENABLED or not signing_status:
        return jsonify({"enabled": False, "message": "integrity API not active"}), 503
    lab_id = _integrity_lab_id()
    status = signing_status(lab_id)
    status["enabled"] = True
    status["signing_available"] = bool(signing_available and signing_available())
    return jsonify(status)


@app.route('/api/integrity/sign/keygen', methods=['POST'])
def api_integrity_sign_keygen():
    if not GETAILAB_INTEGRITY_ENABLED or not generate_keypair:
        return jsonify({"enabled": False, "message": "integrity API not active"}), 503
    lab_id = _integrity_lab_id()
    data = request.get_json(silent=True) or {}
    force = bool(data.get("force")) or request.args.get("force", "").lower() in ("1", "true", "yes")
    result = generate_keypair(lab_id, force=force)
    result["enabled"] = True
    code = 201 if result.get("status") == "generated" else 200
    return jsonify(result), code


@app.route('/api/integrity/sign', methods=['POST'])
def api_integrity_sign():
    """Sign the current saved Merkle root (vault by default)."""
    if not GETAILAB_INTEGRITY_ENABLED or not sign_merkle_tree:
        return jsonify({"enabled": False, "message": "integrity API not active"}), 503
    lab_id = _integrity_lab_id()
    data = request.get_json(silent=True) or {}
    tree_name = data.get("tree_name") or request.args.get("tree", "vault")
    loop_id = data.get("loop_id")
    if loop_id is None and request.args.get("loop", "").isdigit():
        loop_id = int(request.args.get("loop"))
    metadata = {"source": "api_sign", "loop_id": loop_id}
    try:
        record = sign_merkle_tree(lab_id, tree_name=tree_name, metadata=metadata)
        if "error" in record:
            return jsonify({"enabled": True, "ok": False, **record}), 400
        return jsonify({"enabled": True, "ok": True, **record}), 201
    except FileNotFoundError as e:
        return jsonify({"enabled": True, "ok": False, "error": str(e)}), 400
    except Exception as e:
        return jsonify({"enabled": True, "ok": False, "error": str(e)}), 500


@app.route('/api/integrity/sign/verify')
def api_integrity_sign_verify():
    if not GETAILAB_INTEGRITY_ENABLED or not verify_merkle_signature:
        return jsonify({"enabled": False, "message": "integrity API not active"}), 503
    lab_id = _integrity_lab_id()
    tree_name = request.args.get("tree", "vault")
    result = verify_merkle_signature(lab_id, tree_name=tree_name)
    result["enabled"] = True
    result["ok"] = result.get("valid", False)
    return jsonify(result)


@app.route('/api/integrity/sign/attest', methods=['POST'])
def api_integrity_sign_attest():
    """Scan vault Merkle tree and sign root hash (post-archive workflow)."""
    if not GETAILAB_INTEGRITY_ENABLED or not attest_vault:
        return jsonify({"enabled": False, "message": "integrity API not active"}), 503
    lab_id = _integrity_lab_id()
    data = request.get_json(silent=True) or {}
    loop_id = data.get("loop_id")
    if loop_id is None and request.args.get("loop", "").isdigit():
        loop_id = int(request.args.get("loop"))
    sign = data.get("sign", True)
    if "sign" not in data:
        sign = request.args.get("sign", "1").lower() not in ("0", "false", "no")
    result = attest_vault(lab_id, loop_id=loop_id, sign=bool(sign))
    result["enabled"] = True
    result["ok"] = result.get("signed", True) and "error" not in result
    return jsonify(result)


# ============================================================
# JOB TICKETS API (autonomous_core extracted → getailab.tickets)
# ============================================================

@app.route('/api/tickets/summary')
def api_tickets_summary():
    if not GETAILAB_TICKETS_ENABLED or not get_loop_ticket_tracker:
        return jsonify({"enabled": False, "message": "Job tickets not active"})
    tracker = get_loop_ticket_tracker()
    summary = tracker.system.get_daily_summary()
    summary["enabled"] = True
    return jsonify(summary)


@app.route('/api/tickets')
def api_tickets_list():
    if not GETAILAB_TICKETS_ENABLED or not get_loop_ticket_tracker:
        return jsonify({"enabled": False, "tickets": []})
    tracker = get_loop_ticket_tracker()
    assignee = request.args.get("assignee")
    status = request.args.get("status")
    tag = request.args.get("tag")
    limit = min(int(request.args.get("limit", 100)), 500)
    tickets = tracker.system.list_tickets(
        assignee=assignee, status=status, tag=tag, limit=limit
    )
    return jsonify({"enabled": True, "count": len(tickets), "tickets": tickets})


@app.route('/api/tickets/<int:ticket_id>')
def api_ticket_detail(ticket_id):
    if not GETAILAB_TICKETS_ENABLED or not get_loop_ticket_tracker:
        return jsonify({"error": "tickets not active"}), 503
    tracker = get_loop_ticket_tracker()
    ticket = tracker.system.get_ticket(ticket_id)
    if not ticket:
        return jsonify({"error": "ticket not found"}), 404
    ticket["history"] = tracker.system.get_ticket_history(ticket_id)
    return jsonify(ticket)


@app.route('/api/tickets/loop/<int:loop_id>')
def api_tickets_for_loop(loop_id):
    if not GETAILAB_TICKETS_ENABLED or not get_loop_ticket_tracker:
        return jsonify({"enabled": False, "loop_id": loop_id, "tickets": []})
    tracker = get_loop_ticket_tracker()
    return jsonify(tracker.get_loop_summary(loop_id))


@app.route('/api/library/scientist/<scientist>/reference', methods=['GET', 'POST'])
def api_scientist_reference(scientist):
    """
    Beef up a scientist's brain with user-sourced reference material.

    POST JSON: {title, content, url, tags}
    POST multipart: title, tags (comma-separated), file, optional url
    GET: list reference pages (?q=search&limit=20)
    """
    if not GETAILAB_LIBRARY_ENABLED or not add_scientist_reference:
        return jsonify({"ok": False, "error": "library not active"}), 503

    scientist = scientist.lower().strip()
    if not valid_scientist_name(scientist):
        return jsonify({
            "ok": False,
            "error": f"unknown scientist '{scientist}' — use a the example lab squad member (not oracle)",
        }), 400

    if request.method == 'GET':
        query = request.args.get('q', '')
        limit = min(int(request.args.get('limit', 20)), 50)
        try:
            result = get_scientist_references(scientist, query=query, limit=limit)
            return jsonify(result)
        except Exception as e:
            return jsonify({"ok": False, "error": str(e)}), 500

    title = ""
    content = ""
    url = ""
    tags = None
    source_label = "user"

    if request.content_type and "multipart/form-data" in request.content_type:
        title = (request.form.get("title") or "").strip()
        url = (request.form.get("url") or "").strip()
        raw_tags = (request.form.get("tags") or "").strip()
        if raw_tags:
            tags = [t.strip() for t in raw_tags.split(",") if t.strip()]
        upload = request.files.get("file")
        if upload and upload.filename:
            content = upload.read().decode("utf-8", errors="replace")
            source_label = "file"
            if not title:
                title = upload.filename
        else:
            content = (request.form.get("content") or "").strip()
    else:
        data = request.get_json(silent=True) or {}
        title = (data.get("title") or "").strip()
        content = (data.get("content") or "").strip()
        url = (data.get("url") or "").strip()
        raw_tags = data.get("tags")
        if isinstance(raw_tags, list):
            tags = [str(t).strip() for t in raw_tags if str(t).strip()]
        elif isinstance(raw_tags, str) and raw_tags.strip():
            tags = [t.strip() for t in raw_tags.split(",") if t.strip()]

    try:
        result = add_scientist_reference(
            scientist,
            title=title,
            content=content,
            url=url,
            tags=tags,
            source_label=source_label,
        )
        return jsonify(result), 201
    except ValueError as e:
        return jsonify({"ok": False, "error": str(e)}), 400
    except Exception as e:
        return jsonify({"ok": False, "error": str(e)}), 500


if __name__ == '__main__':
    init_db()
    print(f"LAB ONLINE | Lab ID: {ACTIVE_LAB_ID} | isolated vault (no cross-department library bleed)")
    print(f"             | Results DB: {DB_PATH}")
    print(f"             | Artifacts:  {ARTIFACTS_DIR}")
    print(f"             | Agora DB:   {_agora_db_path()}")
    if not is_chimera_lab(ACTIVE_LAB_ID):
        print(f"             | Reports:    {lab_reports_dir(ACTIVE_LAB_ID)}")
    print("Dashboard: /  |  /dashboard  |  /lab")
    print("APIs: /api/stats | /api/loops | /api/library/status | /api/integrity | /api/learner | /api/tickets | /api/chat")
    if GETAILAB_LIBRARY_ENABLED:
        print(f"Library: {GETAILAB_LIB.persistence.root}")
    if GETAILAB_INTEGRITY_ENABLED:
        print("Integrity: /api/integrity/verify | /api/integrity/verify/full | POST /api/integrity/scan | POST /api/integrity/sign/attest")
    if GETAILAB_LEARNING_ENABLED:
        print("Learner: /api/learner/profile | POST /api/learner/interaction | /api/learner/coaching")
    app.run(host='0.0.0.0', port=AGENT_PORT, debug=False)